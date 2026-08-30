from __future__ import annotations

import io

import pytest
from PIL import Image
from pptx import Presentation

from academy.adapters.tools.pptx_writer import (
    SLIDE_DIMENSIONS,
    add_slide,
    create_presentation,
    save_to_bytes,
)


def _image_bytes(width: int, height: int) -> bytes:
    image = Image.new("RGB", (width, height), "white")
    output = io.BytesIO()
    image.save(output, format="PNG")
    return output.getvalue()


@pytest.mark.parametrize(
    ("image_size", "cropped_axis"),
    [
        ((2000, 1000), "horizontal"),
        ((1000, 2000), "vertical"),
    ],
)
def test_cover_fit_uses_real_center_crop_inside_slide(image_size, cropped_axis):
    slide_width, slide_height = SLIDE_DIMENSIONS["16:9"]
    presentation = create_presentation(slide_width, slide_height)

    add_slide(
        presentation,
        _image_bytes(*image_size),
        background_color="white",
        fit_mode="cover",
    )

    presentation = Presentation(io.BytesIO(save_to_bytes(presentation)))
    picture = next(shape for shape in presentation.slides[0].shapes if hasattr(shape, "image"))
    assert (picture.left, picture.top, picture.width, picture.height) == (
        0,
        0,
        slide_width,
        slide_height,
    )
    if cropped_axis == "horizontal":
        assert picture.crop_left == pytest.approx(picture.crop_right)
        assert picture.crop_left > 0
        assert picture.crop_top == 0
        assert picture.crop_bottom == 0
    else:
        assert picture.crop_top == pytest.approx(picture.crop_bottom)
        assert picture.crop_top > 0
        assert picture.crop_left == 0
        assert picture.crop_right == 0
