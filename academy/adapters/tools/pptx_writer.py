# PATH: academy/adapters/tools/pptx_writer.py
# Thin wrapper around python-pptx for PPT file creation.
#
# Isolates the python-pptx dependency from domain logic.

from __future__ import annotations

import io
from typing import Any, Dict, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches


# Slide dimensions in EMU (English Metric Units)
SLIDE_DIMENSIONS: Dict[str, Tuple[int, int]] = {
    "16:9": (Inches(13.333), Inches(7.5)),
    "4:3": (Inches(10), Inches(7.5)),
}

# Background color presets
_BG_COLORS = {
    "black": RGBColor(0, 0, 0),
    "white": RGBColor(255, 255, 255),
    "dark_gray": RGBColor(30, 30, 30),
}


def create_presentation(width: int, height: int) -> Presentation:
    """Create a new Presentation with given slide dimensions.

    Args:
        width: Slide width in EMU.
        height: Slide height in EMU.

    Returns:
        python-pptx Presentation object.
    """
    prs = Presentation()
    prs.slide_width = width
    prs.slide_height = height
    return prs


def _get_blank_layout(prs: Presentation) -> Any:
    """Find the blank slide layout, searching by name then by fewest placeholders."""
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if name in ("blank", "빈 슬라이드", "빈 화면"):
            return layout

    min_ph = None
    best = None
    for layout in prs.slide_layouts:
        ph_count = len(layout.placeholders)
        if min_ph is None or ph_count < min_ph:
            min_ph = ph_count
            best = layout

    return best or prs.slide_layouts[-1]


def _resolve_bg_color(background: str) -> RGBColor:
    """Resolve background color name or hex to RGBColor."""
    color = _BG_COLORS.get(background)
    if color is not None:
        return color

    if background.startswith("#") and len(background) == 7:
        try:
            r = int(background[1:3], 16)
            g = int(background[3:5], 16)
            b = int(background[5:7], 16)
            return RGBColor(r, g, b)
        except ValueError:
            pass

    return _BG_COLORS["black"]


def _fit_image_to_slide(
    img_w: int,
    img_h: int,
    slide_w: int,
    slide_h: int,
    fit_mode: str = "contain",
) -> Tuple[int, int, int, int]:
    """Calculate image placement (left, top, width, height) in EMU."""
    if fit_mode == "stretch":
        return 0, 0, slide_w, slide_h

    if img_w <= 0 or img_h <= 0:
        return 0, 0, slide_w, slide_h

    img_ratio = img_w / img_h
    slide_ratio = slide_w / slide_h

    if fit_mode == "cover":
        if img_ratio > slide_ratio:
            height = slide_h
            width = int(height * img_ratio)
        else:
            width = slide_w
            height = int(width / img_ratio)
    else:  # contain
        if img_ratio > slide_ratio:
            width = slide_w
            height = int(width / img_ratio)
        else:
            height = slide_h
            width = int(height * img_ratio)

    left = (slide_w - width) // 2
    top = (slide_h - height) // 2
    return left, top, width, height


def add_slide(
    prs: Presentation,
    image_bytes: bytes,
    background_color: str = "black",
    fit_mode: str = "contain",
) -> None:
    """Add a slide with an image to the presentation.

    Args:
        prs: Presentation object.
        image_bytes: JPEG or PNG image bytes.
        background_color: Color name or hex string.
        fit_mode: "contain", "cover", or "stretch".
    """
    layout = _get_blank_layout(prs)
    slide = prs.slides.add_slide(layout)

    # Background
    bg_color = _resolve_bg_color(background_color)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Image dimensions
    img = Image.open(io.BytesIO(image_bytes))
    img_w, img_h = img.size

    # Placement
    left, top, width, height = _fit_image_to_slide(
        img_w, img_h, prs.slide_width, prs.slide_height, fit_mode=fit_mode,
    )

    # Add picture
    img_stream = io.BytesIO(image_bytes)
    slide.shapes.add_picture(img_stream, left, top, width, height)


def save_to_bytes(prs: Presentation) -> bytes:
    """Save the presentation to bytes.

    Args:
        prs: Presentation object.

    Returns:
        PPTX file as bytes.
    """
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
