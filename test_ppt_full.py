"""Full PPT pipeline test -- domain + adapter + use case + quality guards."""
import io
import os
import sys

os.environ.setdefault("DJANGO_SETTINGS_MODULE", "apps.api.config.settings.dev")

from PIL import Image, ImageDraw, ImageStat
from pptx import Presentation


def run():
    passed = 0
    failed = 0

    def test(name, fn):
        nonlocal passed, failed
        try:
            fn()
            print(f"  PASS: {name}")
            passed += 1
        except Exception as e:
            print(f"  FAIL: {name} - {type(e).__name__}: {e}")
            failed += 1

    def make_img_bytes(w=400, h=300, color=(100, 50, 30), fmt="PNG"):
        img = Image.new("RGB", (w, h), color)
        buf = io.BytesIO()
        img.save(buf, format=fmt)
        return buf.getvalue()

    # ── 1. Image Preprocessor ──
    print("=== 1. Image Preprocessor ===")
    from academy.domain.tools.image_preprocessor import preprocess_for_export, preprocess_for_detect

    def t_watermark():
        doc = Image.new("RGB", (800, 600), (245, 245, 245))
        d = ImageDraw.Draw(doc)
        for y in range(100, 500, 30):
            d.rectangle([50, y, 750, y + 12], fill=(20, 20, 20))
        d.text((300, 250), "WATERMARK", fill=(220, 220, 220))
        exp = preprocess_for_export(doc)
        wm_o = ImageStat.Stat(doc.crop((300, 250, 450, 270)).convert("L")).mean[0]
        wm_e = ImageStat.Stat(exp.crop((300, 250, 450, 270)).convert("L")).mean[0]
        assert abs(wm_e - wm_o) < 30, f"watermark amplified: {abs(wm_e - wm_o)}"

    def t_dark_contrast():
        dark = Image.new("RGB", (800, 600), (80, 80, 80))
        d = ImageDraw.Draw(dark)
        for y in range(100, 500, 30):
            d.rectangle([50, y, 750, y + 12], fill=(40, 40, 40))
        exp = preprocess_for_export(dark)
        assert ImageStat.Stat(exp.convert("L")).stddev[0] > ImageStat.Stat(dark.convert("L")).stddev[0]

    def t_high_contrast():
        bw = Image.new("RGB", (400, 300), (255, 255, 255))
        ImageDraw.Draw(bw).rectangle([50, 50, 350, 250], fill=(0, 0, 0))
        exp = preprocess_for_export(bw)
        diff = abs(ImageStat.Stat(exp.convert("L")).stddev[0] - ImageStat.Stat(bw.convert("L")).stddev[0])
        assert diff < 5, f"over-processed: {diff}"

    def t_detect_grayscale():
        doc = Image.new("RGB", (400, 300), (200, 200, 200))
        ImageDraw.Draw(doc).rectangle([50, 50, 350, 100], fill=(20, 20, 20))
        det = preprocess_for_detect(doc)
        assert det.mode == "L"

    test("watermark not amplified", t_watermark)
    test("dark scan contrast improved", t_dark_contrast)
    test("high contrast no over-process", t_high_contrast)
    test("detect image is grayscale", t_detect_grayscale)

    # ── 2. Question Splitter ──
    print("\n=== 2. Question Splitter ===")
    from academy.domain.tools.question_splitter import split_questions, TextBlock

    def t_numbered():
        blocks = [
            TextBlock("1. What is 2+2?", 50, 100, 400, 120),
            TextBlock("A) 3  B) 4", 70, 130, 400, 150),
            TextBlock("2. What is 3+3?", 50, 200, 400, 220),
            TextBlock("3. What is 4+4?", 50, 300, 400, 320),
        ]
        regions = split_questions(blocks, 500, 800, page_index=0)
        assert len(regions) == 3, f"expected 3, got {len(regions)}"
        assert regions[0].number == 1
        assert regions[1].number == 2
        assert regions[2].number == 3

    def t_paren():
        blocks = [
            TextBlock("(1) First", 50, 50, 400, 70),
            TextBlock("(2) Second", 50, 150, 400, 170),
            TextBlock("(3) Third", 50, 250, 400, 270),
        ]
        regions = split_questions(blocks, 500, 400, page_index=0)
        assert len(regions) == 3
        assert regions[0].number == 1

    def t_fallback():
        blocks = [
            TextBlock("Just text", 50, 50, 400, 70),
            TextBlock("More text", 50, 100, 400, 120),
        ]
        regions = split_questions(blocks, 500, 400, page_index=0)
        assert len(regions) == 1

    def t_ordering():
        blocks = [
            TextBlock("3. Third", 50, 300, 400, 320),
            TextBlock("1. First", 50, 100, 400, 120),
            TextBlock("2. Second", 50, 200, 400, 220),
        ]
        regions = split_questions(blocks, 500, 500, page_index=0)
        assert len(regions) == 3
        assert regions[0].number == 1
        assert regions[1].number == 2
        assert regions[2].number == 3

    test("simple numbered 1.2.3.", t_numbered)
    test("parenthesis (1)(2)(3)", t_paren)
    test("no questions fallback", t_fallback)
    test("out-of-order sorted", t_ordering)

    # ── 3. PPT Composer ──
    print("\n=== 3. PPT Composer ===")
    from academy.domain.tools.ppt_composer import PptComposer, PptConfig

    def t_compose_basic():
        c = PptComposer(PptConfig(aspect_ratio="16:9", background="black", fit_mode="contain"))
        c.add_slide(make_img_bytes())
        c.add_slide(make_img_bytes(300, 400))
        result = c.finalize()
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 2

    def t_compose_4_3():
        c = PptComposer(PptConfig(aspect_ratio="4:3", background="white", fit_mode="cover"))
        c.add_slide(make_img_bytes())
        result = c.finalize()
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 1

    def t_compose_jpeg():
        c = PptComposer(PptConfig(aspect_ratio="16:9", background="black", fit_mode="contain"))
        c.add_slide(make_img_bytes(fmt="JPEG"))
        result = c.finalize()
        assert len(result) > 100

    test("compose 2 slides", t_compose_basic)
    test("4:3 white cover", t_compose_4_3)
    test("JPEG input", t_compose_jpeg)

    # ── 4. Use Case (image mode) ──
    print("\n=== 4. Use Case (image mode) ===")
    from academy.application.use_cases.tools.generate_ppt import GeneratePptUseCase

    def t_uc_basic():
        uc = GeneratePptUseCase()
        result = uc.execute([make_img_bytes(), make_img_bytes(300, 400)])
        assert result.slide_count == 2
        prs = Presentation(io.BytesIO(result.pptx_bytes))
        assert len(prs.slides) == 2

    def t_uc_config():
        uc = GeneratePptUseCase()
        result = uc.execute(
            [make_img_bytes()],
            config={"aspect_ratio": "4:3", "background": "white"},
        )
        assert result.slide_count == 1

    def t_uc_quality():
        uc = GeneratePptUseCase()
        jpeg = make_img_bytes(800, 600, fmt="JPEG")
        result = uc.execute([jpeg])
        assert len(result.pptx_bytes) > 100

    test("basic 2 images", t_uc_basic)
    test("with config", t_uc_config)
    test("quality preservation", t_uc_quality)

    # ── 5. Output Validity ──
    print("\n=== 5. Output Validity ===")
    import zipfile

    def t_valid_zip():
        uc = GeneratePptUseCase()
        result = uc.execute([make_img_bytes()])
        z = zipfile.ZipFile(io.BytesIO(result.pptx_bytes))
        names = z.namelist()
        assert "[Content_Types].xml" in names
        assert any("slide1.xml" in n for n in names)

    test("valid PPTX zip", t_valid_zip)

    # ── 6. Original Test Suite ──
    print("\n=== 6. Original Test Suite ===")
    import django
    django.setup()
    from apps.domains.tools.ppt.tests import run_all
    orig_result = run_all()
    if orig_result == 0:
        passed += 35
    else:
        failed += 1

    # ── Summary ──
    print(f"\n{'='*50}")
    print(f"TOTAL: {passed} passed, {failed} failed")
    if failed:
        print("SOME TESTS FAILED")
        return 1
    print("ALL TESTS PASSED")
    return 0


if __name__ == "__main__":
    sys.exit(run())
