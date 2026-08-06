from __future__ import annotations

import html
import io
import os
import re
from pathlib import Path
from typing import Any

from apps.domains.tools.problem_review.schema import normalize_report_payload


INK = "14213D"
NAVY = "0B1739"
RED = "D91E3F"
PALE_RED = "FFF0F3"
SLATE = "526079"
LINE = "DCE2EC"
PAPER = "F7F9FC"


def safe_report_filename(value: str, *, suffix: str) -> str:
    stem = re.sub(r"[^0-9A-Za-z가-힣._ -]+", "", str(value or "")).strip(" ._")
    stem = stem[:90] or "문제_리뷰_리포트"
    return f"{stem}_문제리뷰.{suffix}"


def _plain(value: Any) -> str:
    return str(value or "").strip()


def _points_number(value: Any) -> float:
    match = re.search(r"\d+(?:\.\d+)?", _plain(value).replace(",", ""))
    return float(match.group(0)) if match else 0.0


def _exam_structure(report: dict[str, Any]) -> dict[str, Any]:
    questions = report.get("questions") or []
    subjective = [
        item for item in questions
        if re.search(r"서답|서술|주관", _plain(item.get("unit")))
    ]
    objective = [item for item in questions if item not in subjective]
    return {
        "objective_count": len(objective),
        "objective_points": sum(_points_number(item.get("points")) for item in objective),
        "subjective_count": len(subjective),
        "subjective_points": sum(_points_number(item.get("points")) for item in subjective),
    }


def _register_pdf_fonts() -> tuple[str, str]:
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    regular_name = "ProblemReviewRegular"
    bold_name = "ProblemReviewBold"
    try:
        pdfmetrics.getFont(regular_name)
        pdfmetrics.getFont(bold_name)
        return regular_name, bold_name
    except Exception:
        pass

    font_dir = Path(__file__).resolve().parents[2] / "assets" / "omr" / "renderer" / "fonts"
    windows_fonts = Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts"
    regular_candidates = [
        font_dir / "NotoSansKR-Regular.ttf",
        windows_fonts / "malgun.ttf",
        Path("/usr/share/fonts/truetype/nanum/NanumGothic.ttf"),
        Path("/usr/share/fonts/truetype/noto/NotoSansKR-Regular.ttf"),
    ]
    bold_candidates = [
        font_dir / "NotoSansKR-Bold.ttf",
        windows_fonts / "malgunbd.ttf",
        Path("/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf"),
        Path("/usr/share/fonts/truetype/noto/NotoSansKR-Bold.ttf"),
    ]
    registered_regular = False
    for path in regular_candidates:
        if path.is_file():
            try:
                pdfmetrics.registerFont(TTFont(regular_name, str(path)))
                registered_regular = True
                break
            except Exception:
                continue
    if not registered_regular:
        return "Helvetica", "Helvetica-Bold"
    for path in bold_candidates:
        if path.is_file():
            try:
                pdfmetrics.registerFont(TTFont(bold_name, str(path)))
                return regular_name, bold_name
            except Exception:
                continue
    return regular_name, regular_name


def render_problem_review_pdf(payload: dict[str, Any]) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        KeepTogether,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    report = normalize_report_payload(payload, preserve_question_set=False)
    regular, bold = _register_pdf_fonts()
    buffer = io.BytesIO()
    document = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=16 * mm,
        rightMargin=16 * mm,
        topMargin=16 * mm,
        bottomMargin=15 * mm,
        title=_plain(report["metadata"].get("title")) or "문제 리뷰 리포트",
    )
    styles = getSampleStyleSheet()
    body = ParagraphStyle(
        "ReviewBody",
        parent=styles["BodyText"],
        fontName=regular,
        fontSize=9.2,
        leading=14,
        textColor=colors.HexColor(f"#{INK}"),
        wordWrap="CJK",
    )
    small = ParagraphStyle(
        "ReviewSmall",
        parent=body,
        fontSize=7.4,
        leading=10.5,
        textColor=colors.HexColor(f"#{SLATE}"),
    )
    title_style = ParagraphStyle(
        "ReviewTitle",
        parent=body,
        fontName=bold,
        fontSize=25,
        leading=34,
        textColor=colors.white,
    )
    section_style = ParagraphStyle(
        "ReviewSection",
        parent=body,
        fontName=bold,
        fontSize=15,
        leading=20,
        textColor=colors.HexColor(f"#{NAVY}"),
        spaceBefore=6,
        spaceAfter=8,
    )
    card_title = ParagraphStyle(
        "ReviewCardTitle",
        parent=body,
        fontName=bold,
        fontSize=10,
        leading=14,
        textColor=colors.HexColor(f"#{NAVY}"),
    )

    def paragraph(value: Any, style=body):
        text = html.escape(_plain(value)).replace("\n", "<br/>") or "-"
        return Paragraph(text, style)

    def footer(canvas, doc):
        canvas.saveState()
        canvas.setStrokeColor(colors.HexColor(f"#{LINE}"))
        canvas.line(16 * mm, 11 * mm, A4[0] - 16 * mm, 11 * mm)
        canvas.setFont(regular, 7)
        canvas.setFillColor(colors.HexColor(f"#{SLATE}"))
        canvas.drawString(16 * mm, 7 * mm, "업로드 자료 기반 · 선생님 검수형 분석 리포트")
        canvas.drawRightString(A4[0] - 16 * mm, 7 * mm, str(doc.page))
        canvas.restoreState()

    story: list[Any] = []
    meta = report["metadata"]
    cover = Table(
        [[
            Paragraph("PROBLEM REVIEW REPORT", ParagraphStyle(
                "CoverEyebrow", parent=small, fontName=bold, textColor=colors.HexColor(f"#{RED}"),
            )),
        ], [
            paragraph(meta.get("title") or f"{meta.get('school')} {meta.get('exam_name')}", title_style),
        ], [
            paragraph(
                " · ".join(filter(None, [meta.get("school"), meta.get("grade"), meta.get("subject"), meta.get("exam_date")])),
                ParagraphStyle("CoverMeta", parent=body, textColor=colors.HexColor("#D8DEEB")),
            ),
        ]],
        colWidths=[178 * mm],
        rowHeights=[12 * mm, 54 * mm, 20 * mm],
    )
    cover.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor(f"#{NAVY}")),
        ("LEFTPADDING", (0, 0), (-1, -1), 12 * mm),
        ("RIGHTPADDING", (0, 0), (-1, -1), 12 * mm),
        ("TOPPADDING", (0, 0), (-1, -1), 8 * mm),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6 * mm),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LINEBEFORE", (0, 0), (0, -1), 4, colors.HexColor(f"#{RED}")),
    ]))
    story.extend([cover, Spacer(1, 12 * mm)])
    summary = report["summary"]
    structure = _exam_structure(report)
    metric_values = [
        ("문항", str(summary.get("total_questions") or len(report["questions"]))),
        ("선택형", f"{structure['objective_count']}문항"),
        ("서답형", f"{structure['subjective_count']}문항" if structure["subjective_count"] else "확인 필요"),
        ("총점", summary.get("total_points") or "검수 필요"),
    ]
    metric = Table(
        [[paragraph(label, small) for label, _ in metric_values], [paragraph(value, card_title) for _, value in metric_values]],
        colWidths=[44.5 * mm] * 4,
    )
    metric.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor(f"#{PAPER}")),
        ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor(f"#{LINE}")),
        ("INNERGRID", (0, 0), (-1, -1), 0.4, colors.HexColor(f"#{LINE}")),
        ("LEFTPADDING", (0, 0), (-1, -1), 4 * mm),
        ("TOPPADDING", (0, 0), (-1, -1), 2.2 * mm),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 2.2 * mm),
    ]))
    story.extend([
        metric,
        Spacer(1, 3 * mm),
        Table(
            [[paragraph("분석 근거", card_title), paragraph(
                "업로드된 시험지의 문항·배점·자료 구조를 기준으로 분석했습니다. "
                "실제 정답률과 학교 성적 분포가 없는 항목은 추정값으로 확정하지 않습니다.",
                small,
            )]],
            colWidths=[25 * mm, 153 * mm],
            style=TableStyle([
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor(f"#{PAPER}")),
                ("LINEBEFORE", (0, 0), (0, -1), 2.4, colors.HexColor(f"#{RED}")),
                ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor(f"#{LINE}")),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("LEFTPADDING", (0, 0), (-1, -1), 3 * mm),
                ("RIGHTPADDING", (0, 0), (-1, -1), 3 * mm),
                ("TOPPADDING", (0, 0), (-1, -1), 2.5 * mm),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 2.5 * mm),
            ]),
        ),
        Spacer(1, 7 * mm),
        paragraph("시험 한 줄 평", section_style),
        paragraph(summary.get("one_line"), ParagraphStyle("Lead", parent=body, fontName=bold, fontSize=13, leading=20)),
        Spacer(1, 3 * mm),
        paragraph(summary.get("character")),
        Spacer(1, 3 * mm),
        paragraph(summary.get("student_burden"), small),
        PageBreak(),
    ])

    story.append(paragraph("출제 기조와 영역 구성", section_style))
    axes = report["assessment_axes"] or [{"title": "선생님 확인 필요", "description": "출제 기조를 입력해 주세요."}]
    for index, axis in enumerate(axes, start=1):
        table = Table(
            [[paragraph(f"AXIS {index}", ParagraphStyle("Axis", parent=small, fontName=bold, textColor=colors.HexColor(f"#{RED}"))),
              paragraph(axis.get("title"), card_title)],
             ["", paragraph(axis.get("description"))]],
            colWidths=[23 * mm, 155 * mm],
        )
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), colors.white),
            ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor(f"#{LINE}")),
            ("SPAN", (0, 1), (0, 1)),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 3.5 * mm),
            ("RIGHTPADDING", (0, 0), (-1, -1), 3.5 * mm),
            ("TOPPADDING", (0, 0), (-1, -1), 2.4 * mm),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 2.4 * mm),
        ]))
        story.extend([table, Spacer(1, 2.5 * mm)])

    if report["domains"]:
        story.extend([Spacer(1, 4 * mm), paragraph("영역별 분석", section_style)])
        domain_rows = [[paragraph("영역", card_title), paragraph("문항", card_title), paragraph("배점/비중", card_title), paragraph("해석", card_title)]]
        for item in report["domains"]:
            numbers = ", ".join(item.get("question_numbers") or [])
            domain_rows.append([
                paragraph(item.get("name")),
                paragraph(numbers, small),
                paragraph(" / ".join(filter(None, [item.get("points"), item.get("ratio")])), small),
                paragraph(item.get("insight"), small),
            ])
        domain_table = Table(domain_rows, colWidths=[35 * mm, 32 * mm, 27 * mm, 84 * mm], repeatRows=1)
        domain_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{NAVY}")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor(f"#{LINE}")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(f"#{PAPER}")]),
            ("LEFTPADDING", (0, 0), (-1, -1), 2.2 * mm),
            ("RIGHTPADDING", (0, 0), (-1, -1), 2.2 * mm),
            ("TOPPADDING", (0, 0), (-1, -1), 2 * mm),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 2 * mm),
        ]))
        story.append(domain_table)

    distributions = report["difficulty"].get("distribution") or []
    if distributions:
        story.extend([Spacer(1, 4 * mm), paragraph("난이도 분포", section_style)])
        difficulty_rows = [[
            paragraph("난이도", card_title),
            paragraph("문항", card_title),
            paragraph("배점", card_title),
            paragraph("검수 메모", card_title),
        ]]
        for item in distributions:
            difficulty_rows.append([
                paragraph(item.get("label"), small),
                paragraph(", ".join(item.get("question_numbers") or []), small),
                paragraph(item.get("points"), small),
                paragraph(item.get("note"), small),
            ])
        difficulty_table = Table(
            difficulty_rows,
            colWidths=[24 * mm, 42 * mm, 25 * mm, 87 * mm],
            repeatRows=1,
        )
        difficulty_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{NAVY}")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor(f"#{LINE}")),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(f"#{PAPER}")]),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 2.2 * mm),
            ("RIGHTPADDING", (0, 0), (-1, -1), 2.2 * mm),
            ("TOPPADDING", (0, 0), (-1, -1), 2 * mm),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 2 * mm),
        ]))
        story.append(difficulty_table)
        if report["difficulty"].get("grade_estimate_note"):
            story.extend([
                Spacer(1, 2.5 * mm),
                paragraph(report["difficulty"].get("grade_estimate_note"), small),
            ])

    story.extend([PageBreak(), paragraph("전 문항 리뷰", section_style)])
    question_rows = [[
        paragraph("번호", card_title), paragraph("단원", card_title), paragraph("배점", card_title),
        paragraph("난이도", card_title), paragraph("핵심 포인트", card_title), paragraph("주요 함정", card_title),
    ]]
    for item in report["questions"]:
        question_rows.append([
            paragraph(item.get("number"), small),
            paragraph(item.get("unit"), small),
            paragraph(item.get("points"), small),
            paragraph(item.get("difficulty"), small),
            paragraph(item.get("key_point") or item.get("review_note"), small),
            paragraph(item.get("trap"), small),
        ])
    questions_table = Table(
        question_rows,
        colWidths=[11 * mm, 28 * mm, 16 * mm, 19 * mm, 58 * mm, 46 * mm],
        repeatRows=1,
    )
    questions_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{NAVY}")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor(f"#{LINE}")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(f"#{PAPER}")]),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("ALIGN", (0, 0), (3, -1), "CENTER"),
        ("LEFTPADDING", (0, 0), (-1, -1), 1.6 * mm),
        ("RIGHTPADDING", (0, 0), (-1, -1), 1.6 * mm),
        ("TOPPADDING", (0, 0), (-1, -1), 1.8 * mm),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 1.8 * mm),
    ]))
    story.append(questions_table)

    story.extend([PageBreak(), paragraph("핵심 변별 문항", section_style)])
    for item in report["key_items"] or [{"rank": 1, "title": "선생님 확인 필요"}]:
        content = Table(
            [[paragraph(f"RANK {item.get('rank')}", ParagraphStyle("Rank", parent=small, fontName=bold, textColor=colors.HexColor(f"#{RED}"))),
              paragraph(item.get("title"), card_title)],
             [paragraph("왜 어려운가", small), paragraph(item.get("reason"))],
             [paragraph("무너지는 지점", small), paragraph(item.get("collapse_point"))],
             [paragraph("학습 처방", small), paragraph(item.get("prescription"))]],
            colWidths=[31 * mm, 147 * mm],
        )
        content.setStyle(TableStyle([
            ("BOX", (0, 0), (-1, -1), 0.7, colors.HexColor(f"#{LINE}")),
            ("BACKGROUND", (0, 0), (0, -1), colors.HexColor(f"#{PALE_RED}")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 3 * mm),
            ("RIGHTPADDING", (0, 0), (-1, -1), 3 * mm),
            ("TOPPADDING", (0, 0), (-1, -1), 2.4 * mm),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 2.4 * mm),
        ]))
        story.extend([KeepTogether(content), Spacer(1, 4 * mm)])

    story.extend([paragraph("학생이 무너지는 패턴", section_style)])
    for index, item in enumerate(report["failure_patterns"], start=1):
        story.extend([
            paragraph(f"{index}. {item.get('title')}", card_title),
            paragraph(
                f"증상: {item.get('symptom')}\n원인: {item.get('cause')}\n처방: {item.get('prescription')}",
                body,
            ),
            Spacer(1, 3 * mm),
        ])

    guidance = report["parent_guidance"]
    if guidance.get("avoid") or guidance.get("recommended"):
        story.extend([Spacer(1, 3 * mm), paragraph("학부모 설명 가이드", section_style)])
        guidance_table = Table(
            [[paragraph("피할 표현", card_title), paragraph("권장 설명", card_title)], [
                paragraph("\n".join(f"- {item}" for item in guidance.get("avoid") or []), small),
                paragraph("\n".join(f"- {item}" for item in guidance.get("recommended") or []), small),
            ]],
            colWidths=[89 * mm, 89 * mm],
        )
        guidance_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (0, -1), colors.HexColor(f"#{PALE_RED}")),
            ("BACKGROUND", (1, 0), (1, -1), colors.HexColor(f"#{PAPER}")),
            ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor(f"#{LINE}")),
            ("INNERGRID", (0, 0), (-1, -1), 0.4, colors.HexColor(f"#{LINE}")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 3 * mm),
            ("RIGHTPADDING", (0, 0), (-1, -1), 3 * mm),
            ("TOPPADDING", (0, 0), (-1, -1), 2.5 * mm),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 2.5 * mm),
        ]))
        story.append(KeepTogether(guidance_table))

    story.extend([PageBreak(), paragraph("다음 시험까지의 결론", section_style)])
    story.append(paragraph(report["conclusion"].get("headline"), ParagraphStyle(
        "Conclusion", parent=body, fontName=bold, fontSize=17, leading=25, textColor=colors.HexColor(f"#{NAVY}"),
    )))
    story.append(Spacer(1, 5 * mm))
    for index, action in enumerate(report["conclusion"].get("actions") or [], start=1):
        story.append(paragraph(f"{index}. {action}", body))
        story.append(Spacer(1, 2 * mm))
    if report["warnings"]:
        story.extend([Spacer(1, 8 * mm), paragraph("검수 전 확인", section_style)])
        for warning in report["warnings"]:
            story.append(paragraph(f"- {warning}", small))

    document.build(story, onFirstPage=footer, onLaterPages=footer)
    return buffer.getvalue()


def render_problem_review_pptx(payload: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    report = normalize_report_payload(payload, preserve_question_set=False)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rgb(value: str) -> RGBColor:
        return RGBColor.from_string(value)

    def rect(slide, x, y, w, h, *, fill: str, line: str | None = None):
        from pptx.enum.shapes import MSO_SHAPE

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
        shape.line.color.rgb = rgb(line or fill)
        return shape

    def text_box(
        slide,
        value: Any,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        size: float = 18,
        color: str = INK,
        bold: bool = False,
        align=PP_ALIGN.LEFT,
        font: str = "Noto Sans KR",
        margin: float = 0.05,
        valign=MSO_ANCHOR.TOP,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.margin_left = Inches(margin)
        frame.margin_right = Inches(margin)
        frame.margin_top = Inches(margin)
        frame.margin_bottom = Inches(margin)
        frame.vertical_anchor = valign
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = _plain(value)
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return box

    def chrome(slide, eyebrow: str, title: str, number: int):
        rect(slide, 0, 0, 0.12, 7.5, fill=RED)
        text_box(slide, eyebrow.upper(), 0.35, 0.26, 5.8, 0.28, size=9, color=RED, bold=True)
        text_box(slide, title, 0.35, 0.58, 12.2, 0.58, size=25, color=NAVY, bold=True)
        text_box(slide, f"{number} / {{total}}", 11.8, 7.1, 1.0, 0.2, size=8, color=SLATE, align=PP_ALIGN.RIGHT)
        rect(slide, 0.35, 1.24, 12.45, 0.015, fill=LINE)

    meta = report["metadata"]
    slide = prs.slides.add_slide(blank)
    rect(slide, 0, 0, 13.333333, 7.5, fill=NAVY)
    rect(slide, 0, 0, 0.18, 7.5, fill=RED)
    text_box(slide, "PROBLEM REVIEW REPORT", 0.6, 0.72, 5.0, 0.3, size=11, color=RED, bold=True)
    text_box(slide, meta.get("title") or f"{meta.get('school')} {meta.get('exam_name')}", 0.6, 1.42, 11.5, 1.5, size=34, color="FFFFFF", bold=True)
    text_box(slide, " · ".join(filter(None, [meta.get("school"), meta.get("grade"), meta.get("subject"), meta.get("exam_date")])), 0.62, 3.15, 10.6, 0.45, size=15, color="D8DEEB")
    text_box(slide, "업로드 자료 기반 · 선생님 검수형 분석 리포트", 0.62, 6.72, 4.8, 0.25, size=10, color="AAB5CB")

    slide = prs.slides.add_slide(blank)
    chrome(slide, "EXAM OVERVIEW", "시험 개요와 한 줄 평", 2)
    summary = report["summary"]
    structure = _exam_structure(report)
    metrics = [
        ("문항", str(summary.get("total_questions") or len(report["questions"]))),
        ("선택형", f"{structure['objective_count']}문항"),
        ("서답형", f"{structure['subjective_count']}문항" if structure["subjective_count"] else "확인 필요"),
        ("총점", summary.get("total_points") or "검수 필요"),
    ]
    for index, (label, value) in enumerate(metrics):
        x = 0.35 + index * 2.05
        rect(slide, x, 1.55, 1.85, 1.05, fill=PAPER, line=LINE)
        text_box(slide, label, x + 0.12, 1.72, 1.6, 0.2, size=9, color=SLATE)
        text_box(slide, value, x + 0.12, 2.02, 1.6, 0.32, size=19, color=NAVY, bold=True)
    rect(slide, 0.35, 2.95, 12.45, 1.12, fill=NAVY)
    text_box(slide, summary.get("one_line") or "한 줄 평을 입력해 주세요.", 0.7, 3.2, 11.7, 0.62, size=20, color="FFFFFF", bold=True, valign=MSO_ANCHOR.MIDDLE)
    text_box(slide, summary.get("character"), 0.48, 4.35, 12.0, 0.78, size=14, color=INK)
    rect(slide, 0.48, 5.36, 12.0, 0.68, fill=PALE_RED, line="F6C4CE")
    text_box(slide, summary.get("student_burden") or "학생 부담 요인을 검수해 주세요.", 0.72, 5.51, 11.5, 0.34, size=12, color=INK)
    rect(slide, 0.48, 6.18, 12.0, 0.48, fill=PAPER, line=LINE)
    text_box(slide, "분석 근거 · 업로드 시험지의 문항·배점·자료 구조 기준 / 실제 정답률·학교 성적 분포는 포함하지 않음", 0.72, 6.31, 11.5, 0.22, size=9, color=SLATE)

    slide = prs.slides.add_slide(blank)
    chrome(slide, "ASSESSMENT AXES", "이 시험이 확인하는 역량", 3)
    axes = report["assessment_axes"][:6]
    for index, axis in enumerate(axes):
        col = index % 2
        row = index // 2
        x = 0.4 + col * 6.25
        y = 1.55 + row * 1.62
        rect(slide, x, y, 5.95, 1.35, fill=PAPER, line=LINE)
        text_box(slide, f"AXIS {index + 1}", x + 0.18, y + 0.14, 1.05, 0.2, size=9, color=RED, bold=True)
        text_box(slide, axis.get("title"), x + 1.25, y + 0.11, 4.45, 0.32, size=15, color=NAVY, bold=True)
        text_box(slide, axis.get("description"), x + 0.18, y + 0.53, 5.55, 0.64, size=11, color=SLATE)

    if report["domains"]:
        slide = prs.slides.add_slide(blank)
        chrome(slide, "DOMAIN BREAKDOWN", "영역별 출제 비중", 4)
        rows = min(7, len(report["domains"])) + 1
        table_height = min(4.95, max(1.5, rows * 0.68))
        table_shape = slide.shapes.add_table(
            rows,
            4,
            Inches(0.45),
            Inches(1.58),
            Inches(12.35),
            Inches(table_height),
        )
        table = table_shape.table
        widths = [2.3, 2.25, 1.7, 6.1]
        for index, width in enumerate(widths):
            table.columns[index].width = Inches(width)
        values = [["영역", "문항", "배점/비중", "해석"]] + [
            [item.get("name"), ", ".join(item.get("question_numbers") or []), " / ".join(filter(None, [item.get("points"), item.get("ratio")])), item.get("insight")]
            for item in report["domains"][:7]
        ]
        for row_index, row in enumerate(values):
            for col_index, value in enumerate(row):
                cell = table.cell(row_index, col_index)
                cell.text = _plain(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(NAVY if row_index == 0 else ("FFFFFF" if row_index % 2 else PAPER))
                cell.margin_left = Inches(0.09)
                cell.margin_right = Inches(0.09)
                cell.margin_top = Inches(0.06)
                cell.margin_bottom = Inches(0.06)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = "Noto Sans KR"
                        run.font.size = Pt(10 if row_index else 11)
                        run.font.bold = row_index == 0
                        run.font.color.rgb = rgb("FFFFFF" if row_index == 0 else INK)

    distributions = report["difficulty"].get("distribution") or []
    if distributions:
        slide = prs.slides.add_slide(blank)
        chrome(slide, "DIFFICULTY MAP", "난이도 분포와 해석 주의", len(prs.slides))
        rows = min(6, len(distributions)) + 1
        table_height = min(4.35, max(1.6, rows * 0.66))
        table_shape = slide.shapes.add_table(
            rows,
            4,
            Inches(0.48),
            Inches(1.55),
            Inches(12.25),
            Inches(table_height),
        )
        table = table_shape.table
        for index, width in enumerate([1.35, 2.55, 1.4, 6.95]):
            table.columns[index].width = Inches(width)
        values = [["난이도", "문항", "배점", "검수 메모"]] + [
            [item.get("label"), ", ".join(item.get("question_numbers") or []), item.get("points"), item.get("note")]
            for item in distributions[:6]
        ]
        for row_index, row in enumerate(values):
            for col_index, value in enumerate(row):
                cell = table.cell(row_index, col_index)
                cell.text = _plain(value) or "-"
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(NAVY if row_index == 0 else ("FFFFFF" if row_index % 2 else PAPER))
                cell.margin_left = Inches(0.09)
                cell.margin_right = Inches(0.09)
                cell.margin_top = Inches(0.05)
                cell.margin_bottom = Inches(0.05)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER if col_index < 3 else PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.name = "Noto Sans KR"
                        run.font.size = Pt(10)
                        run.font.bold = row_index == 0
                        run.font.color.rgb = rgb("FFFFFF" if row_index == 0 else INK)
        if report["difficulty"].get("grade_estimate_note"):
            rect(slide, 0.48, 6.18, 12.25, 0.62, fill=PALE_RED, line="F6C4CE")
            text_box(slide, report["difficulty"].get("grade_estimate_note"), 0.72, 6.32, 11.75, 0.32, size=10.5, color=INK)

    question_slide_start = len(prs.slides) + 1
    question_groups = [report["questions"][start:start + 6] for start in range(0, len(report["questions"]), 6)] or [[]]
    for group_index, group in enumerate(question_groups):
        slide = prs.slides.add_slide(blank)
        first = group[0].get("number") if group else "-"
        last = group[-1].get("number") if group else "-"
        chrome(slide, "ITEM ANALYSIS", f"문항 {first}~{last}번 분석표", question_slide_start + group_index)
        rows = len(group) + 1
        table_height = min(5.45, max(1.75, rows * 0.72))
        table_shape = slide.shapes.add_table(
            rows,
            6,
            Inches(0.35),
            Inches(1.5),
            Inches(12.6),
            Inches(table_height),
        )
        table = table_shape.table
        for index, width in enumerate([0.6, 1.55, 0.75, 0.9, 4.65, 4.15]):
            table.columns[index].width = Inches(width)
        values = [["번호", "단원", "배점", "난이도", "핵심 포인트", "주요 함정"]] + [
            [item.get("number"), item.get("unit"), item.get("points"), item.get("difficulty"), item.get("key_point") or item.get("review_note"), item.get("trap")]
            for item in group
        ]
        for row_index, row in enumerate(values):
            for col_index, value in enumerate(row):
                cell = table.cell(row_index, col_index)
                cell.text = _plain(value) or "-"
                cell.fill.solid()
                highlight = row_index > 0 and row[3] in {"상", "최상"}
                cell.fill.fore_color.rgb = rgb(NAVY if row_index == 0 else (PALE_RED if highlight else ("FFFFFF" if row_index % 2 else PAPER)))
                cell.margin_left = Inches(0.07)
                cell.margin_right = Inches(0.07)
                cell.margin_top = Inches(0.04)
                cell.margin_bottom = Inches(0.04)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER if col_index < 4 else PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.name = "Noto Sans KR"
                        run.font.size = Pt(9.5)
                        run.font.bold = row_index == 0
                        run.font.color.rgb = rgb("FFFFFF" if row_index == 0 else INK)

    for key_index, item in enumerate(report["key_items"][:5], start=1):
        slide = prs.slides.add_slide(blank)
        chrome(slide, f"KILLER REVIEW #{key_index}", item.get("title") or "핵심 변별 문항", len(prs.slides))
        rect(slide, 0.42, 1.55, 2.2, 4.95, fill=PALE_RED, line="F6C4CE")
        text_box(slide, f"RANK {item.get('rank') or key_index}", 0.65, 1.92, 1.75, 0.25, size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)
        text_box(slide, ", ".join(item.get("question_numbers") or []) or "검수 필요", 0.65, 2.45, 1.75, 0.75, size=26, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        blocks = [
            ("왜 어려운가", item.get("reason")),
            ("학생이 무너지는 지점", item.get("collapse_point")),
            ("다음 시험 처방", item.get("prescription")),
        ]
        for block_index, (label, value) in enumerate(blocks):
            y = 1.58 + block_index * 1.63
            rect(slide, 2.95, y, 9.75, 1.38, fill="FFFFFF", line=LINE)
            text_box(slide, label, 3.18, y + 0.15, 2.2, 0.25, size=11, color=RED, bold=True)
            text_box(slide, value, 3.18, y + 0.53, 9.1, 0.64, size=13, color=INK)

    guidance = report["parent_guidance"]
    if guidance.get("avoid") or guidance.get("recommended"):
        slide = prs.slides.add_slide(blank)
        chrome(slide, "TEACHER COMMUNICATION", "학생·학부모 설명 가이드", len(prs.slides))
        rect(slide, 0.45, 1.5, 5.95, 4.85, fill=PALE_RED, line="F6C4CE")
        text_box(slide, "피할 표현", 0.72, 1.74, 2.0, 0.25, size=11, color=RED, bold=True)
        avoid_text = "\n".join(f"• {item}" for item in guidance.get("avoid") or []) or "선생님 확인 필요"
        text_box(slide, avoid_text, 0.72, 2.16, 5.3, 3.72, size=15, color=INK)
        rect(slide, 6.72, 1.5, 6.05, 4.85, fill=PAPER, line=LINE)
        text_box(slide, "권장 설명", 6.99, 1.74, 2.0, 0.25, size=11, color=NAVY, bold=True)
        recommended_text = "\n".join(f"• {item}" for item in guidance.get("recommended") or []) or "선생님 확인 필요"
        text_box(slide, recommended_text, 6.99, 2.16, 5.42, 3.72, size=15, color=INK)

    failure_patterns = report["failure_patterns"]
    for page_start in range(0, len(failure_patterns), 4):
        slide = prs.slides.add_slide(blank)
        page_patterns = failure_patterns[page_start:page_start + 4]
        chrome(slide, "LEARNING DIAGNOSIS", "학생이 무너지는 패턴", len(prs.slides))
        for offset, item in enumerate(page_patterns):
            index = page_start + offset
            col = offset % 2
            row = offset // 2
            x = 0.45 + col * 6.25
            y = 1.5 + row * 2.62
            rect(slide, x, y, 6.05, 2.3, fill="FFFFFF", line=LINE)
            text_box(slide, f"PATTERN {index + 1}", x + 0.22, y + 0.24, 1.05, 0.22, size=8.5, color=RED, bold=True)
            text_box(slide, item.get("title"), x + 1.28, y + 0.18, 4.45, 0.48, size=13, color=NAVY, bold=True)
            pattern_body = (
                f"증상 · {_plain(item.get('symptom'))}\n"
                f"원인 · {_plain(item.get('cause'))}\n"
                f"처방 · {_plain(item.get('prescription'))}"
            )
            text_box(slide, pattern_body, x + 0.22, y + 0.78, 5.58, 1.28, size=9.8, color=SLATE)

    slide = prs.slides.add_slide(blank)
    chrome(slide, "FINAL TAKEAWAY", "다음 시험까지 무엇을 다질 것인가", len(prs.slides))
    rect(slide, 0.45, 1.55, 12.25, 1.15, fill=NAVY)
    text_box(slide, report["conclusion"].get("headline") or summary.get("one_line"), 0.78, 1.84, 11.6, 0.56, size=20, color="FFFFFF", bold=True, valign=MSO_ANCHOR.MIDDLE)
    actions = report["conclusion"].get("actions") or []
    for index, action in enumerate(actions[:6]):
        col = index % 2
        row = index // 2
        x = 0.48 + col * 6.18
        y = 3.1 + row * 1.18
        rect(slide, x, y, 5.88, 0.92, fill=PAPER, line=LINE)
        rect(slide, x + 0.14, y + 0.18, 0.5, 0.5, fill=RED)
        text_box(slide, index + 1, x + 0.14, y + 0.2, 0.5, 0.3, size=11, color="FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        text_box(slide, action, x + 0.82, y + 0.19, 4.75, 0.5, size=13, color=INK, bold=True)

    total = len(prs.slides)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "{total}" in run.text:
                        run.text = run.text.replace("{total}", str(total))

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def render_problem_review_report(payload: dict[str, Any], *, output_format: str) -> tuple[bytes, str, str]:
    report = normalize_report_payload(payload, preserve_question_set=False)
    export_meta = dict(payload.get("_export_meta", {})) if isinstance(payload, dict) and isinstance(payload.get("_export_meta"), dict) else {}
    verification = payload.get("verification") if isinstance(payload, dict) else None
    if isinstance(verification, dict) and verification.get("status") == "verified":
        export_meta.setdefault("review_completed_at", verification.get("verified_at"))
        export_meta.setdefault("source_fingerprint", verification.get("report_fingerprint"))
    render_payload = {
        **report,
        "_export_meta": export_meta,
    }
    title = report.get("metadata", {}).get("title") or "문제 리뷰 리포트"
    if output_format == "pdf":
        return (
            render_problem_review_pdf(render_payload),
            safe_report_filename(title, suffix="pdf"),
            "application/pdf",
        )
    if output_format == "pptx":
        return (
            render_problem_review_pptx(render_payload),
            safe_report_filename(title, suffix="pptx"),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    raise ValueError("PDF 또는 PPTX만 선택할 수 있습니다.")


# EXAM SPECTRUM is the current product renderer.  The legacy functions above are
# intentionally kept in this module for historical output compatibility while
# all new exports resolve these names to the evidence-led design system.
from apps.domains.tools.problem_review.spectrum_renderers import (  # noqa: E402
    render_problem_review_pdf as _render_problem_review_pdf_spectrum,
    render_problem_review_pptx as _render_problem_review_pptx_spectrum,
)

render_problem_review_pdf = _render_problem_review_pdf_spectrum
render_problem_review_pptx = _render_problem_review_pptx_spectrum
