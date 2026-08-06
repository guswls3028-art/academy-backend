from __future__ import annotations

import html
import io
import math
import os
import re
from collections import Counter
from pathlib import Path
from typing import Any, Iterable

from apps.domains.tools.problem_review.schema import normalize_report_payload


DEEP_INK = "09162F"
PLASMA_BLUE = "37B7FF"
SIGNAL_CORAL = "FF526F"
ION_AMBER = "F4B746"
LAB_PAPER = "F5F7FB"
CARBON = "172033"
MIST = "DDE5F0"
MUTED = "65738B"
WHITE = "FFFFFF"
ACTION_COLORS = {
    "확인": "7CD4FF",
    "해석": PLASMA_BLUE,
    "계산": ION_AMBER,
    "서술": SIGNAL_CORAL,
    "복합": "9B8CFF",
    "검수 필요": "9AA6B8",
}
DIFFICULTY_LEVEL = {"검수 필요": 0, "하": 1, "중": 2, "중상": 3, "상": 4, "최상": 5}


def _plain(value: Any) -> str:
    text = str(value or "").strip()
    return re.sub(r"DNA\s*양", "DNA 양", text)


def _number(value: Any) -> float | None:
    match = re.search(r"-?\d+(?:\.\d+)?", _plain(value).replace(",", ""))
    if not match:
        return None
    try:
        return float(match.group(0))
    except ValueError:
        return None


def _fmt_number(value: float) -> str:
    return str(int(value)) if float(value).is_integer() else f"{value:.1f}".rstrip("0").rstrip(".")


def _export_meta(payload: dict[str, Any]) -> dict[str, Any]:
    value = payload.get("_export_meta") if isinstance(payload, dict) else None
    return value if isinstance(value, dict) else {}


def _data_contract(payload: dict[str, Any]) -> tuple[dict[str, Any], dict[str, Any]]:
    report = normalize_report_payload(payload, preserve_question_set=False)
    questions = report["questions"]
    point_values = [_number(item.get("points")) for item in questions]
    all_points_known = bool(questions) and all(value is not None for value in point_values)
    total_points = sum(value or 0 for value in point_values) if all_points_known else None
    numbers = [item["number"] for item in questions]
    duplicate_numbers = len(numbers) != len(set(numbers))
    limitations = [
        "실제 정답률·학교 성적 분포는 포함하지 않음",
        *(report.get("warnings") or []),
    ]
    if not all_points_known:
        limitations.append("일부 문항의 배점이 확인되지 않아 배점 합계와 비중은 표시하지 않습니다.")
    if duplicate_numbers:
        limitations.append("표시 문항 번호가 중복되어 단원별 배점 교차검산은 선생님 확인이 필요합니다.")
    if not report["difficulty"].get("grade_estimate_note"):
        limitations.append("실제 점수 분포가 없어 등급컷과 성취 비율은 제시하지 않습니다.")
    for field, label in (("school", "학교"), ("subject", "과목"), ("exam_name", "시험명")):
        if not report["metadata"].get(field):
            limitations.append(f"{label} 정보가 없어 내보내기 전에 입력이 필요합니다.")
    metrics = {
        "question_count": len(questions),
        "total_points": total_points,
        "all_points_known": all_points_known,
        "duplicate_numbers": duplicate_numbers,
        "limitations": list(dict.fromkeys(item for item in limitations if item)),
    }
    return report, metrics


def _domain_for_question(report: dict[str, Any], question: dict[str, Any]) -> str:
    number = str(question.get("number"))
    for domain in report.get("domains") or []:
        if number in {str(value) for value in domain.get("question_numbers") or []}:
            return _plain(domain.get("name")) or _plain(question.get("unit")) or "미분류"
    return _plain(question.get("unit")) or "미분류"


def _balanced_chunks(items: list[Any], target: int) -> list[list[Any]]:
    if not items:
        return [[]]
    group_count = max(1, math.ceil(len(items) / max(1, target)))
    base, extra = divmod(len(items), group_count)
    groups: list[list[Any]] = []
    cursor = 0
    for index in range(group_count):
        size = base + (1 if index < extra else 0)
        groups.append(items[cursor: cursor + size])
        cursor += size
    return groups


def _split_text(value: Any, limit: int) -> list[str]:
    text = re.sub(r"\s+", " ", _plain(value))
    if not text:
        return [""]
    chunks: list[str] = []
    while len(text) > limit:
        break_at = text.rfind(" ", 0, limit + 1)
        if break_at < max(12, limit // 3):
            break_at = limit
        chunks.append(text[:break_at].strip())
        text = text[break_at:].strip()
    chunks.append(text)
    return chunks


def _ledger_rows(report: dict[str, Any]) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for question in report.get("questions") or []:
        key_parts = _split_text(question.get("key_point") or question.get("review_note"), 130)
        trap_parts = _split_text(question.get("trap"), 110)
        count = max(len(key_parts), len(trap_parts))
        for index in range(count):
            rows.append({
                "number": str(question.get("number")) if index == 0 else "계속",
                "points": _plain(question.get("points")) if index == 0 else "",
                "domain": _domain_for_question(report, question) if index == 0 else "",
                "action": _plain(question.get("thinking_action")) if index == 0 else "",
                "difficulty": _plain(question.get("difficulty")) if index == 0 else "",
                "key": key_parts[index] if index < len(key_parts) else "",
                "trap": trap_parts[index] if index < len(trap_parts) else "",
            })
    return rows


def _enriched_key_items(report: dict[str, Any]) -> list[dict[str, Any]]:
    by_number = {str(item.get("number")): item for item in report.get("questions") or []}
    conclusion_actions = list(report.get("conclusion", {}).get("actions") or [])
    output: list[dict[str, Any]] = []
    for raw in (report.get("key_items") or [])[:3]:
        item = dict(raw)
        related = [by_number[number] for number in map(str, item.get("question_numbers") or []) if number in by_number]
        evidence_parts = [item.get("evidence") or item.get("reason")]
        evidence_parts.extend(question.get("key_point") for question in related)
        item["evidence"] = "\n".join(dict.fromkeys(_plain(value) for value in evidence_parts if _plain(value)))
        branches = list(item.get("collapse_branches") or [])
        for question in related:
            branches.extend([question.get("trap"), question.get("validity"), question.get("review_note")])
        branches = list(dict.fromkeys(_plain(value) for value in branches if _plain(value)))
        while len(branches) < 3:
            branches.append(f"분기 {len(branches) + 1} · 선생님 검수 필요")
        item["collapse_branches"] = branches[:3]
        steps = list(item.get("recovery_steps") or [])
        steps.extend([item.get("prescription"), *conclusion_actions])
        steps = list(dict.fromkeys(_plain(value) for value in steps if _plain(value)))
        while len(steps) < 4:
            steps.append(f"복구 {len(steps) + 1}단계 · 선생님 검수 필요")
        item["recovery_steps"] = steps[:4]
        item["learning_point"] = _plain(item.get("learning_point") or item.get("prescription"))
        output.append(item)
    return output


def _xray_items(report: dict[str, Any]) -> list[dict[str, Any]]:
    """Return the same evidence-complete top three X-rays for every output format."""
    output = _enriched_key_items(report)
    candidates = sorted(
        report.get("questions") or [],
        key=lambda item: (DIFFICULTY_LEVEL.get(item.get("difficulty"), 0), _number(item.get("points")) or 0),
        reverse=True,
    )
    actions = list(report.get("conclusion", {}).get("actions") or [])
    while len(output) < min(3, len(candidates)):
        question = candidates[len(output)]
        branches = list(dict.fromkeys(
            _plain(value)
            for value in (question.get("trap"), question.get("validity"), question.get("review_note"))
            if _plain(value)
        ))
        while len(branches) < 3:
            branches.append(f"분기 {len(branches) + 1} · 선생님 검수 필요")
        recovery = list(dict.fromkeys(
            _plain(value)
            for value in (question.get("review_note"), *actions)
            if _plain(value)
        ))
        while len(recovery) < 4:
            recovery.append(f"복구 {len(recovery) + 1}단계 · 선생님 검수 필요")
        output.append({
            "rank": len(output) + 1,
            "title": f"{question.get('number')}번 · {_domain_for_question(report, question)}",
            "question_numbers": [str(question.get("number"))],
            "evidence": question.get("key_point"),
            "collapse_branches": branches[:3],
            "recovery_steps": recovery[:4],
            "learning_point": question.get("key_point"),
        })
    return output[:3]


def _pptx_text_style(run, *, font: str, size: float, color, bold: bool = False) -> None:
    from pptx.util import Pt

    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def render_problem_review_pptx(payload: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    report, metrics = _data_contract(payload)
    export_meta = _export_meta(payload)
    version = int(export_meta.get("report_version") or 0)
    fingerprint = _plain(export_meta.get("source_fingerprint"))
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rgb(value: str) -> RGBColor:
        return RGBColor.from_string(value)

    def shape(slide, x, y, w, h, *, fill: str, line: str | None = None, radius: bool = False):
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        item = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
        item.fill.solid()
        item.fill.fore_color.rgb = rgb(fill)
        item.line.color.rgb = rgb(line or fill)
        return item

    def line(slide, x1, y1, x2, y2, *, color: str = MIST, width: float = 1.0):
        item = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        item.line.color.rgb = rgb(color)
        item.line.width = Pt(width)
        return item

    def text(
        slide,
        value: Any,
        x: float,
        y: float,
        w: float,
        h: float,
        *,
        size: float = 14,
        color: str = CARBON,
        bold: bool = False,
        align=PP_ALIGN.LEFT,
        font: str = "Pretendard",
        valign=MSO_ANCHOR.TOP,
        margin: float = 0.02,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(margin)
        frame.margin_right = Inches(margin)
        frame.margin_top = Inches(margin)
        frame.margin_bottom = Inches(margin)
        frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = _plain(value)
        _pptx_text_style(run, font=font, size=size, color=rgb(color), bold=bold)
        return box

    def add_notes(slide, summary: str) -> None:
        meta = report["metadata"]
        sources = [
            f"- {meta.get('school') or '학교 미입력'} · {meta.get('subject') or '과목 미입력'} · {meta.get('exam_name') or '시험명 미입력'}",
            f"- 검수본 v{version or '-'} · snapshot {fingerprint[:12] or '미지정'}",
            "- 선생님이 저장한 문제 리뷰 검수본",
        ]
        slide.notes_slide.notes_text_frame.text = "\n".join([
            _plain(summary),
            "",
            "[Sources]",
            *sources,
        ])

    slides_for_numbering = []

    def base_slide(*, rail: str, title: str, subtitle: str = ""):
        slide = prs.slides.add_slide(blank)
        shape(slide, 0, 0, 13.333333, 7.5, fill=LAB_PAPER)
        shape(slide, 0.28, 0.36, 0.06, 6.62, fill=PLASMA_BLUE)
        shape(slide, 0.19, 0.66, 0.24, 0.06, fill=SIGNAL_CORAL)
        text(slide, rail, 0.52, 0.35, 2.6, 0.24, size=9, color=SIGNAL_CORAL, bold=True)
        text(slide, title, 0.52, 0.7, 12.1, 0.52, size=25, color=DEEP_INK, bold=True)
        if subtitle:
            text(slide, subtitle, 0.54, 1.2, 11.8, 0.3, size=10, color=MUTED)
        line(slide, 0.52, 1.55, 12.82, 1.55, color=MIST, width=0.8)
        text(slide, "EXAM SPECTRUM", 0.53, 7.12, 2.0, 0.18, size=7.5, color=MUTED, bold=True)
        slides_for_numbering.append(slide)
        return slide

    meta = report["metadata"]
    summary = report["summary"]
    questions = report["questions"]

    # 1. Observation cover
    slide = prs.slides.add_slide(blank)
    shape(slide, 0, 0, 13.333333, 7.5, fill=DEEP_INK)
    shape(slide, 0.28, 0, 0.08, 7.5, fill=PLASMA_BLUE)
    shape(slide, 0.28, 0.78, 0.08, 1.35, fill=SIGNAL_CORAL)
    text(slide, "OBSERVATION RECORD / EXAM SPECTRUM", 0.7, 0.68, 5.8, 0.28, size=10, color=PLASMA_BLUE, bold=True)
    text(slide, meta.get("school") or "학교 정보 확인 필요", 0.7, 1.28, 10.8, 0.72, size=35, color=WHITE, bold=True)
    text(slide, " · ".join(filter(None, [meta.get("subject"), meta.get("grade"), meta.get("exam_name")])), 0.72, 2.18, 10.5, 0.36, size=16, color="C8D4E6")
    text(slide, summary.get("one_line") or "시험의 한 문장 신호를 검수해 주세요.", 0.72, 3.05, 11.7, 1.2, size=25, color=WHITE, bold=True)
    line(slide, 0.72, 4.72, 12.35, 4.72, color="334761", width=1.0)
    spectrum_width = 11.6 / max(1, len(questions))
    for index, question in enumerate(questions):
        action = question.get("thinking_action") or "검수 필요"
        level = DIFFICULTY_LEVEL.get(question.get("difficulty"), 0)
        x = 0.73 + index * spectrum_width
        h = 0.16 + level * 0.055
        shape(slide, x, 5.08 - h, max(0.025, spectrum_width * 0.68), h, fill=ACTION_COLORS.get(action, "9AA6B8"))
    metric_value = f"{metrics['question_count']}문항"
    if metrics["total_points"] is not None:
        metric_value += f" · {_fmt_number(metrics['total_points'])}점"
    text(slide, metric_value, 0.72, 5.46, 3.2, 0.42, size=19, color=ION_AMBER, bold=True, font="Aptos Mono")
    text(slide, f"검수본 v{version or '-'} · {fingerprint[:12] or 'snapshot 미지정'}", 0.72, 6.74, 5.0, 0.22, size=8.5, color="8292AA", font="Aptos Mono")
    slides_for_numbering.append(slide)
    add_notes(slide, "학교, 시험, 한 줄 평과 전체 문항 스펙트럼을 여는 관측 표지입니다.")

    # 2. Three-minute briefing
    slide = base_slide(rail="3-MINUTE BRIEF", title="시험을 세 문장으로 읽습니다", subtitle="무엇을 물었고 · 어디서 무너졌고 · 다음에 무엇을 바꿀지")
    lanes = [
        ("SIGNAL", "무엇을 물었나", summary.get("character") or summary.get("one_line"), PLASMA_BLUE),
        ("BREAKPOINT", "어디서 무너졌나", summary.get("student_burden") or "학생 부담 요인을 검수해 주세요.", SIGNAL_CORAL),
        ("NEXT", "무엇을 바꿀까", report["conclusion"].get("headline") or "다음 시험 행동을 검수해 주세요.", ION_AMBER),
    ]
    for index, (code, label, value, color) in enumerate(lanes):
        y = 1.82 + index * 1.33
        shape(slide, 0.53, y, 0.08, 0.96, fill=color)
        text(slide, code, 0.77, y, 1.15, 0.22, size=8.5, color=color, bold=True, font="Aptos Mono")
        text(slide, label, 1.95, y - 0.02, 2.15, 0.34, size=17, color=DEEP_INK, bold=True)
        text(slide, value, 4.05, y - 0.03, 8.45, 0.75, size=13, color=CARBON)
        line(slide, 0.77, y + 1.04, 12.62, y + 1.04, color=MIST, width=0.6)
    limitation = metrics["limitations"][0] if metrics["limitations"] else "정답·배점·문항 수는 저장된 검수본을 기준으로 교차검산했습니다."
    text(slide, f"DATA LIMIT · {limitation}", 0.77, 6.12, 11.75, 0.42, size=9.5, color=MUTED)
    add_notes(slide, "시험 성격, 붕괴 지점, 다음 행동과 데이터 한계를 3분 안에 설명하는 브리핑입니다.")

    # 3. Evaluation DNA matrix
    slide = base_slide(rail="EVALUATION DNA", title="단원과 사고행동이 만나는 지점", subtitle="문항 수가 아니라 어떤 행동을 요구했는지까지 함께 봅니다.")
    domains = [item.get("name") or "미분류" for item in report.get("domains") or []]
    if not domains:
        domains = list(dict.fromkeys(_domain_for_question(report, item) for item in questions))[:8]
    domains = domains[:7]
    actions = ["확인", "해석", "계산", "서술", "복합", "검수 필요"]
    matrix = Counter((_domain_for_question(report, item), item.get("thinking_action") or "검수 필요") for item in questions)
    x0, y0 = 3.05, 1.9
    cell_w, cell_h = 1.48, min(0.62, 4.55 / max(1, len(domains)))
    for col, action in enumerate(actions):
        text(slide, action, x0 + col * cell_w, 1.66, cell_w, 0.22, size=9, color=ACTION_COLORS[action], bold=True, align=PP_ALIGN.CENTER)
    for row, domain in enumerate(domains):
        y = y0 + row * cell_h
        text(slide, domain, 0.58, y + 0.13, 2.26, 0.25, size=10.5, color=DEEP_INK, bold=True)
        line(slide, 0.55, y + cell_h, 12.35, y + cell_h, color=MIST, width=0.5)
        for col, action in enumerate(actions):
            count = matrix[(domain, action)]
            if count:
                diameter = 0.18 + min(0.24, count * 0.055)
                marker = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(x0 + col * cell_w + (cell_w - diameter) / 2),
                    Inches(y + (cell_h - diameter) / 2),
                    Inches(diameter),
                    Inches(diameter),
                )
                marker.fill.solid(); marker.fill.fore_color.rgb = rgb(ACTION_COLORS[action]); marker.line.color.rgb = rgb(ACTION_COLORS[action])
                text(slide, count, x0 + col * cell_w, y + 0.11, cell_w, 0.28, size=8, color=DEEP_INK, bold=True, align=PP_ALIGN.CENTER, font="Aptos Mono")
    add_notes(slide, "단원별 문항을 확인, 해석, 계산, 서술, 복합 사고행동으로 교차 분류한 평가 DNA입니다.")

    # 4. Terrain
    slide = base_slide(rail="TEST TERRAIN", title="출제 지형은 배점과 난도를 겹쳐 읽습니다", subtitle="배점이 확인되지 않은 자료는 문항 수만 표시하고 비중을 만들지 않습니다.")
    domain_rows = report.get("domains") or [{"name": domain} for domain in domains]
    for index, domain in enumerate(domain_rows[:7]):
        y = 1.86 + index * 0.68
        name = domain.get("name") or "미분류"
        domain_questions = [item for item in questions if _domain_for_question(report, item) == name]
        count = len(domain_questions)
        total = sum((_number(item.get("points")) or 0) for item in domain_questions)
        known = bool(domain_questions) and all(_number(item.get("points")) is not None for item in domain_questions)
        width = 7.0 * count / max(1, len(questions))
        shape(slide, 3.0, y + 0.08, max(0.18, width), 0.25, fill=PLASMA_BLUE)
        for q_index, item in enumerate(domain_questions):
            level = DIFFICULTY_LEVEL.get(item.get("difficulty"), 0)
            x = 3.02 + (q_index + 0.5) * max(0.18, width) / max(1, count)
            marker_color = SIGNAL_CORAL if level >= 4 else ION_AMBER if level == 3 else WHITE
            shape(slide, x, y + 0.02, 0.035, 0.37, fill=marker_color)
        text(slide, name, 0.58, y, 2.15, 0.3, size=11, color=DEEP_INK, bold=True)
        value = f"{count}문항" + (f" · {_fmt_number(total)}점" if known else " · 배점 검수 필요")
        text(slide, value, 10.28, y, 2.15, 0.3, size=10, color=MUTED, align=PP_ALIGN.RIGHT, font="Aptos Mono")
        line(slide, 0.58, y + 0.53, 12.45, y + 0.53, color=MIST, width=0.45)
    add_notes(slide, "단원별 문항 수와 확인된 배점 위에 상위 난도 문항의 위치를 겹쳐 본 출제 지형입니다.")

    # 5. Signature spectrum
    slide = base_slide(rail="EXAM SPECTRUM", title="시험 전체를 한 줄의 증거로 연결합니다", subtitle="색은 사고행동, 높이는 난도, 숫자는 실제 문항 번호입니다.")
    groups = _balanced_chunks(questions, 13)
    row_h = 1.75
    for row_index, group in enumerate(groups[:3]):
        y = 2.15 + row_index * row_h
        line(slide, 0.82, y + 0.55, 12.35, y + 0.55, color="AEBBCB", width=1.2)
        step = 11.45 / max(1, len(group))
        for index, item in enumerate(group):
            x = 0.88 + index * step
            level = DIFFICULTY_LEVEL.get(item.get("difficulty"), 0)
            action = item.get("thinking_action") or "검수 필요"
            height = 0.22 + level * 0.085
            shape(slide, x, y + 0.55 - height, max(0.07, step * 0.48), height, fill=ACTION_COLORS.get(action, "9AA6B8"))
            text(slide, item.get("number"), x - 0.06, y + 0.68, max(0.28, step * 0.6), 0.2, size=8.2, color=DEEP_INK, bold=True, align=PP_ALIGN.CENTER, font="Aptos Mono")
            text(slide, _plain(item.get("points")) or "-", x - 0.08, y + 0.95, max(0.32, step * 0.7), 0.18, size=6.8, color=MUTED, align=PP_ALIGN.CENTER, font="Aptos Mono")
    legend_x = 0.75
    for action in actions:
        shape(slide, legend_x, 6.3, 0.12, 0.12, fill=ACTION_COLORS[action])
        text(slide, action, legend_x + 0.17, 6.24, 0.74, 0.2, size=7.8, color=MUTED)
        legend_x += 1.25
    add_notes(slide, "모든 문항을 순서, 난도, 배점, 사고행동으로 연결한 이 리포트의 고유 EXAM SPECTRUM입니다.")

    # 6+. Evidence ledger, dynamically paginated without truncation.
    ledger_groups = _balanced_chunks(_ledger_rows(report), 9)
    for page_index, group in enumerate(ledger_groups, start=1):
        slide = base_slide(rail="EVIDENCE LEDGER", title=f"전 문항 증거 원장 · {page_index}/{len(ledger_groups)}", subtitle="번호 · 배점 · 단원 · 사고행동 · 난도 · 핵심 · 함정을 한 줄에서 대조합니다.")
        rows = len(group) + 1
        table_shape = slide.shapes.add_table(rows, 7, Inches(0.52), Inches(1.78), Inches(12.28), Inches(5.05))
        table = table_shape.table
        for index, width in enumerate([0.62, 0.66, 1.5, 0.72, 0.72, 4.0, 4.06]):
            table.columns[index].width = Inches(width)
        values = [["번호", "배점", "단원", "행동", "난도", "핵심 증거", "무너지는 함정"]] + [
            [row["number"], row["points"], row["domain"], row["action"], row["difficulty"], row["key"], row["trap"]]
            for row in group
        ]
        for r_index, row in enumerate(values):
            for c_index, value in enumerate(row):
                cell = table.cell(r_index, c_index)
                cell.text = _plain(value) or "-"
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(DEEP_INK if r_index == 0 else (WHITE if r_index % 2 else LAB_PAPER))
                cell.margin_left = Inches(0.055); cell.margin_right = Inches(0.055)
                cell.margin_top = Inches(0.035); cell.margin_bottom = Inches(0.035)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER if c_index < 5 else PP_ALIGN.LEFT
                    for run in paragraph.runs:
                        run.font.name = "Aptos Mono" if c_index < 2 else "Pretendard"
                        run.font.size = Pt(7.2 if r_index else 8.2)
                        run.font.bold = r_index == 0
                        run.font.color.rgb = rgb(WHITE if r_index == 0 else CARBON)
        add_notes(slide, f"전 문항 원장의 {page_index}번째 구간입니다. 표시된 문장을 줄이지 않고 다음 구간으로 분할합니다.")

    # Condition accumulation map
    slide = base_slide(rail="CONDITION MAP", title="조건이 누적될수록 오류는 연쇄됩니다", subtitle="복합·계산·서술 문항 가운데 난도와 배점을 함께 확인합니다.")
    candidates = sorted(
        questions,
        key=lambda item: (DIFFICULTY_LEVEL.get(item.get("difficulty"), 0), _number(item.get("points")) or 0),
        reverse=True,
    )[:6]
    for index, item in enumerate(candidates):
        y = 1.83 + index * 0.78
        color = ACTION_COLORS.get(item.get("thinking_action") or "검수 필요", "9AA6B8")
        text(slide, item.get("number"), 0.64, y, 0.54, 0.34, size=17, color=color, bold=True, align=PP_ALIGN.CENTER, font="Aptos Mono")
        shape(slide, 1.38, y + 0.09, 1.1, 0.12, fill=color)
        text(slide, _domain_for_question(report, item), 2.7, y - 0.02, 2.0, 0.31, size=11, color=DEEP_INK, bold=True)
        text(slide, item.get("key_point"), 4.72, y - 0.03, 3.7, 0.48, size=9.5, color=CARBON)
        text(slide, item.get("trap") or "함정 검수 필요", 8.62, y - 0.03, 3.75, 0.48, size=9.5, color=MUTED)
        line(slide, 0.66, y + 0.57, 12.44, y + 0.57, color=MIST, width=0.5)
    add_notes(slide, "난도와 배점이 높은 문항에서 핵심 조건과 함정이 어떻게 연결되는지 보여 주는 조건 누적 지도입니다.")

    # Three question X-rays.
    key_items = _xray_items(report)
    for xray_index, item in enumerate(key_items, start=1):
        slide = base_slide(rail=f"QUESTION X-RAY {xray_index}", title=item.get("title") or "핵심 변별 문항", subtitle=f"문항 {', '.join(item.get('question_numbers') or []) or '번호 검수 필요'} · 증거 → 붕괴 분기 → 복구 순서")
        text(slide, "EVIDENCE", 0.68, 1.88, 1.4, 0.22, size=9, color=PLASMA_BLUE, bold=True, font="Aptos Mono")
        text(slide, item.get("evidence") or item.get("reason"), 0.68, 2.26, 3.35, 2.85, size=14, color=DEEP_INK, bold=True)
        line(slide, 4.18, 1.9, 4.18, 5.92, color=PLASMA_BLUE, width=2)
        text(slide, "BREAK BRANCHES", 4.52, 1.88, 2.0, 0.22, size=9, color=SIGNAL_CORAL, bold=True, font="Aptos Mono")
        branches = item.get("collapse_branches") or [item.get("collapse_point")]
        for index, branch in enumerate(branches[:3]):
            y = 2.32 + index * 1.15
            shape(slide, 4.53, y + 0.03, 0.08, 0.76, fill=SIGNAL_CORAL)
            text(slide, branch, 4.82, y, 3.1, 0.88, size=11, color=CARBON)
        line(slide, 8.15, 1.9, 8.15, 5.92, color=ION_AMBER, width=2)
        text(slide, "RECOVERY STEPS", 8.48, 1.88, 2.1, 0.22, size=9, color=ION_AMBER, bold=True, font="Aptos Mono")
        steps = item.get("recovery_steps") or [item.get("prescription")]
        for index, step_value in enumerate(steps[:4]):
            y = 2.23 + index * 0.86
            text(slide, f"{index + 1:02d}", 8.48, y, 0.52, 0.28, size=11, color=ION_AMBER, bold=True, font="Aptos Mono")
            text(slide, step_value, 9.05, y - 0.02, 3.2, 0.62, size=10.5, color=CARBON)
        text(slide, f"LEARNING SIGNAL · {item.get('learning_point') or item.get('prescription')}", 4.52, 6.13, 7.72, 0.42, size=9.5, color=DEEP_INK, bold=True)
        add_notes(slide, "핵심 문항의 근거, 붕괴 분기, 네 단계 복구와 학습 신호를 연결한 X-ray입니다.")

    # Error genome
    slide = base_slide(rail="ERROR GENOME", title="표면 실수를 근본 원인과 처방에 연결합니다", subtitle="‘실수했다’로 끝내지 않고 반복되는 오류 구조를 기록합니다.")
    patterns = (report.get("failure_patterns") or [])[:4]
    for index, item in enumerate(patterns):
        y = 1.83 + index * 1.18
        text(slide, f"G{index + 1:02d}", 0.62, y, 0.62, 0.28, size=11, color=SIGNAL_CORAL, bold=True, font="Aptos Mono")
        text(slide, item.get("title"), 1.38, y - 0.02, 2.15, 0.34, size=13, color=DEEP_INK, bold=True)
        text(slide, item.get("symptom"), 3.65, y - 0.02, 2.55, 0.76, size=9.5, color=CARBON)
        line(slide, 6.3, y + 0.25, 6.72, y + 0.25, color=SIGNAL_CORAL, width=1.5)
        text(slide, item.get("cause"), 6.88, y - 0.02, 2.48, 0.76, size=9.5, color=CARBON)
        line(slide, 9.46, y + 0.25, 9.88, y + 0.25, color=ION_AMBER, width=1.5)
        text(slide, item.get("prescription"), 10.02, y - 0.02, 2.35, 0.76, size=9.5, color=CARBON)
        line(slide, 0.62, y + 0.92, 12.38, y + 0.92, color=MIST, width=0.45)
    add_notes(slide, "학생 오류를 증상, 원인, 처방 순서로 연결해 반복 패턴을 확인합니다.")

    # Recovery protocol + achievement bands
    slide = base_slide(rail="RECOVERY PROTOCOL", title="72시간 · 2주 · 다음 시험으로 행동을 나눕니다", subtitle="성취 구간은 실제 점수가 아니라 관찰 신호와 학습 행동으로 설명합니다.")
    protocol = report.get("recovery_protocol") or {}
    columns = [
        ("72H", "72시간 안", protocol.get("within_72_hours") or report["conclusion"].get("actions") or [], PLASMA_BLUE),
        ("2W", "2주 안", protocol.get("within_two_weeks") or [], ION_AMBER),
        ("NEXT", "다음 시험", protocol.get("next_exam") or [], SIGNAL_CORAL),
    ]
    for index, (code, label, items, color) in enumerate(columns):
        x = 0.62 + index * 4.04
        shape(slide, x, 1.86, 0.08, 3.02, fill=color)
        text(slide, code, x + 0.26, 1.86, 0.72, 0.22, size=10, color=color, bold=True, font="Aptos Mono")
        text(slide, label, x + 0.26, 2.2, 2.5, 0.34, size=17, color=DEEP_INK, bold=True)
        body = "\n".join(f"• {_plain(item)}" for item in items[:5]) or "• 선생님 행동 계획 검수 필요"
        text(slide, body, x + 0.26, 2.77, 3.28, 1.75, size=10.5, color=CARBON)
    bands = report.get("achievement_bands") or []
    text(slide, "ACHIEVEMENT SIGNALS", 0.62, 5.24, 2.2, 0.22, size=9, color=MUTED, bold=True, font="Aptos Mono")
    band_text = "  /  ".join(f"{item.get('label')}: {item.get('signal')} → {item.get('prescription')}" for item in bands[:3])
    text(slide, band_text or "점수 분포 대신 풀이 과정에서 관찰할 신호를 선생님이 입력해 주세요.", 0.62, 5.62, 11.72, 0.86, size=9.5, color=CARBON)
    add_notes(slide, "복구 행동을 72시간, 2주, 다음 시험으로 나누고 관찰 가능한 성취 신호를 연결합니다.")

    # Parent conversation memo
    slide = base_slide(rail="PARENT MEMO", title="점수보다 먼저 확인할 질문", subtitle="피할 말과 권장 질문을 나란히 두어 상담 문장으로 바로 사용합니다.")
    guidance = report.get("parent_guidance") or {}
    shape(slide, 0.62, 1.87, 0.09, 4.54, fill=SIGNAL_CORAL)
    text(slide, "피할 말", 0.94, 1.88, 2.1, 0.34, size=18, color=SIGNAL_CORAL, bold=True)
    text(slide, "\n".join(f"— {_plain(item)}" for item in guidance.get("avoid") or []) or "— 선생님 확인 필요", 0.94, 2.45, 4.85, 3.48, size=13, color=CARBON)
    line(slide, 6.28, 1.87, 6.28, 6.4, color=MIST, width=1.2)
    shape(slide, 6.7, 1.87, 0.09, 4.54, fill=PLASMA_BLUE)
    text(slide, "함께 확인할 질문", 7.02, 1.88, 3.2, 0.34, size=18, color=PLASMA_BLUE, bold=True)
    text(slide, "\n".join(f"— {_plain(item)}" for item in guidance.get("recommended") or []) or "— 선생님 확인 필요", 7.02, 2.45, 5.15, 3.48, size=13, color=CARBON)
    add_notes(slide, "학부모 상담에서 피할 표현과 확인할 질문을 점수 중심 표현과 분리합니다.")

    # Next signal
    slide = prs.slides.add_slide(blank)
    shape(slide, 0, 0, 13.333333, 7.5, fill=DEEP_INK)
    shape(slide, 0.3, 0.62, 0.08, 5.88, fill=PLASMA_BLUE)
    text(slide, "NEXT SIGNAL", 0.78, 0.75, 2.1, 0.3, size=10, color=SIGNAL_CORAL, bold=True, font="Aptos Mono")
    text(slide, report["conclusion"].get("headline") or summary.get("one_line"), 0.78, 1.48, 11.5, 1.48, size=29, color=WHITE, bold=True)
    actions_text = "\n".join(f"{index + 1:02d}  {_plain(item)}" for index, item in enumerate(report["conclusion"].get("actions") or []))
    text(slide, actions_text or "01  다음 시험 행동을 검수해 주세요.", 0.8, 3.38, 9.65, 1.92, size=14, color="C8D4E6", font="Aptos Mono")
    text(slide, f"{meta.get('school') or '학교 미입력'} · {meta.get('subject') or '과목 미입력'} · {meta.get('exam_name') or '시험명 미입력'}", 0.8, 6.15, 8.5, 0.28, size=11, color=PLASMA_BLUE, bold=True)
    text(slide, f"v{version or '-'} / {fingerprint[:12] or 'snapshot 미지정'}", 9.55, 6.15, 2.75, 0.28, size=9.5, color="8292AA", align=PP_ALIGN.RIGHT, font="Aptos Mono")
    slides_for_numbering.append(slide)
    add_notes(slide, "다음 시험에서 바꿀 한 가지와 저장된 검수본의 버전, 스냅샷을 고정해 공유합니다.")

    total = len(prs.slides)
    for index, slide in enumerate(prs.slides, start=1):
        if index not in {1, total}:
            text(slide, f"{index:02d} / {total:02d}", 11.62, 7.08, 0.92, 0.18, size=7.5, color=MUTED, align=PP_ALIGN.RIGHT, font="Aptos Mono")
        text(slide, f"v{version or '-'} · {fingerprint[:8] or 'snapshot'}", 9.38, 7.08, 1.98, 0.18, size=7, color=MUTED, align=PP_ALIGN.RIGHT, font="Aptos Mono")

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _register_pdf_fonts() -> tuple[str, str, str]:
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    regular, bold = "SpectrumRegular", "SpectrumBold"
    try:
        pdfmetrics.getFont(regular); pdfmetrics.getFont(bold)
        return regular, bold, regular
    except Exception:
        pass
    font_dir = Path(__file__).resolve().parents[2] / "assets" / "omr" / "renderer" / "fonts"
    windows = Path(os.environ.get("WINDIR", r"C:\Windows")) / "Fonts"
    pairs = [
        (font_dir / "NotoSansKR-Regular.ttf", font_dir / "NotoSansKR-Bold.ttf"),
        (windows / "malgun.ttf", windows / "malgunbd.ttf"),
        (Path("/usr/share/fonts/truetype/nanum/NanumGothic.ttf"), Path("/usr/share/fonts/truetype/nanum/NanumGothicBold.ttf")),
    ]
    for regular_path, bold_path in pairs:
        if regular_path.is_file() and bold_path.is_file():
            pdfmetrics.registerFont(TTFont(regular, str(regular_path)))
            pdfmetrics.registerFont(TTFont(bold, str(bold_path)))
            return regular, bold, regular
    return "Helvetica", "Helvetica-Bold", "Courier"


def render_problem_review_pdf(payload: dict[str, Any]) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_RIGHT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import KeepTogether, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    report, metrics = _data_contract(payload)
    export_meta = _export_meta(payload)
    version = int(export_meta.get("report_version") or 0)
    fingerprint = _plain(export_meta.get("source_fingerprint"))
    regular, bold, mono = _register_pdf_fonts()
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=15 * mm,
        rightMargin=15 * mm,
        topMargin=16 * mm,
        bottomMargin=16 * mm,
        title=_plain(report["metadata"].get("title")) or "EXAM SPECTRUM",
    )
    styles = getSampleStyleSheet()
    body = ParagraphStyle("SpectrumBody", parent=styles["BodyText"], fontName=regular, fontSize=8.8, leading=13.2, textColor=colors.HexColor(f"#{CARBON}"), wordWrap="CJK")
    small = ParagraphStyle("SpectrumSmall", parent=body, fontSize=7.2, leading=10.2, textColor=colors.HexColor(f"#{MUTED}"))
    title = ParagraphStyle("SpectrumTitle", parent=body, fontName=bold, fontSize=23, leading=31, textColor=colors.white)
    section = ParagraphStyle("SpectrumSection", parent=body, fontName=bold, fontSize=15, leading=20, textColor=colors.HexColor(f"#{DEEP_INK}"), spaceBefore=5, spaceAfter=7)
    h3 = ParagraphStyle("SpectrumH3", parent=body, fontName=bold, fontSize=10.5, leading=14, textColor=colors.HexColor(f"#{DEEP_INK}"))
    code = ParagraphStyle("SpectrumCode", parent=small, fontName=mono, fontSize=7, leading=9, textColor=colors.HexColor(f"#{SIGNAL_CORAL}"))

    def p(value: Any, style=body):
        return Paragraph(html.escape(_plain(value)).replace("\n", "<br/>") or "-", style)

    def footer(canvas, current_doc):
        canvas.saveState()
        canvas.setStrokeColor(colors.HexColor(f"#{MIST}")); canvas.line(15 * mm, 11 * mm, A4[0] - 15 * mm, 11 * mm)
        canvas.setFont(mono, 6.5); canvas.setFillColor(colors.HexColor(f"#{MUTED}"))
        canvas.drawString(15 * mm, 7 * mm, f"EXAM SPECTRUM · v{version or '-'} · {fingerprint[:10] or 'snapshot'}")
        canvas.drawRightString(A4[0] - 15 * mm, 7 * mm, str(current_doc.page))
        canvas.restoreState()

    story: list[Any] = []
    meta, summary = report["metadata"], report["summary"]
    cover = Table([
        [p("OBSERVATION RECORD / EXAM SPECTRUM", code)],
        [p(meta.get("school") or "학교 정보 확인 필요", title)],
        [p(" · ".join(filter(None, [meta.get("subject"), meta.get("grade"), meta.get("exam_name")])), ParagraphStyle("CoverMeta", parent=body, fontSize=11, textColor=colors.HexColor("#C8D4E6")))],
        [p(summary.get("one_line") or "시험의 한 문장 신호를 검수해 주세요.", ParagraphStyle("CoverSignal", parent=title, fontSize=16, leading=23))],
        [p(f"{metrics['question_count']}문항" + (f" · {_fmt_number(metrics['total_points'])}점" if metrics["total_points"] is not None else " · 배점 검수 필요"), ParagraphStyle("CoverMetric", parent=body, fontName=mono, fontSize=13, textColor=colors.HexColor(f"#{ION_AMBER}")))],
    ], colWidths=[180 * mm], rowHeights=[14 * mm, 35 * mm, 15 * mm, 55 * mm, 24 * mm])
    cover.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor(f"#{DEEP_INK}")),
        ("LINEBEFORE", (0, 0), (0, -1), 5, colors.HexColor(f"#{PLASMA_BLUE}")),
        ("LEFTPADDING", (0, 0), (-1, -1), 11 * mm), ("RIGHTPADDING", (0, 0), (-1, -1), 8 * mm),
        ("TOPPADDING", (0, 0), (-1, -1), 5 * mm), ("BOTTOMPADDING", (0, 0), (-1, -1), 4 * mm),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    story.extend([cover, PageBreak()])

    story.extend([p("3-MINUTE BRIEF", code), p("시험을 세 문장으로 읽습니다", section)])
    briefing = [
        ("무엇을 물었나", summary.get("character") or summary.get("one_line"), PLASMA_BLUE),
        ("어디서 무너졌나", summary.get("student_burden"), SIGNAL_CORAL),
        ("무엇을 바꿀까", report["conclusion"].get("headline"), ION_AMBER),
    ]
    for label, value, color in briefing:
        row = Table([[p(label, h3), p(value)]], colWidths=[38 * mm, 142 * mm])
        row.setStyle(TableStyle([
            ("LINEBEFORE", (0, 0), (0, 0), 3, colors.HexColor(f"#{color}")),
            ("LINEBELOW", (0, 0), (-1, -1), 0.5, colors.HexColor(f"#{MIST}")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 3 * mm), ("RIGHTPADDING", (0, 0), (-1, -1), 2 * mm),
            ("TOPPADDING", (0, 0), (-1, -1), 3 * mm), ("BOTTOMPADDING", (0, 0), (-1, -1), 3 * mm),
        ]))
        story.extend([row, Spacer(1, 2 * mm)])
    if metrics["limitations"]:
        story.extend([Spacer(1, 4 * mm), p("DATA LIMIT", code), p("\n".join(f"- {item}" for item in metrics["limitations"]), small)])

    story.extend([PageBreak(), p("EVALUATION DNA", code), p("출제 기조와 단원 지도", section)])
    axis_rows = [[p("평가 축", h3), p("관측 내용", h3)]]
    for axis in report.get("assessment_axes") or []:
        axis_rows.append([p(axis.get("title"), h3), p(axis.get("description"), small)])
    axis_table = Table(axis_rows, colWidths=[42 * mm, 138 * mm], repeatRows=1)
    axis_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{DEEP_INK}")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("LINEBEFORE", (0, 1), (0, -1), 2.5, colors.HexColor(f"#{PLASMA_BLUE}")),
        ("LINEBELOW", (0, 0), (-1, -1), 0.35, colors.HexColor(f"#{MIST}")), ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 2.5 * mm), ("RIGHTPADDING", (0, 0), (-1, -1), 2.5 * mm),
        ("TOPPADDING", (0, 0), (-1, -1), 2 * mm), ("BOTTOMPADDING", (0, 0), (-1, -1), 2 * mm),
    ]))
    story.extend([axis_table, Spacer(1, 6 * mm), p("TEST TERRAIN", code), p("단원별 문항 · 배점 · 비중", section)])
    domain_rows = [[p("단원", h3), p("문항", h3), p("배점/비중", h3), p("해석", h3)]]
    for domain in report.get("domains") or []:
        name = domain.get("name") or "미분류"
        items = [item for item in report["questions"] if _domain_for_question(report, item) == name]
        values = [_number(item.get("points")) for item in items]
        points = f"{_fmt_number(sum(value or 0 for value in values))}점" if items and all(value is not None for value in values) else "검수 필요"
        domain_rows.append([p(name), p(", ".join(str(item.get("number")) for item in items), small), p(points, small), p(domain.get("insight"), small)])
    domain_table = Table(domain_rows, colWidths=[34 * mm, 36 * mm, 25 * mm, 85 * mm], repeatRows=1)
    domain_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{DEEP_INK}")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(f"#{LAB_PAPER}")]),
        ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor(f"#{MIST}")), ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("LEFTPADDING", (0, 0), (-1, -1), 2 * mm), ("RIGHTPADDING", (0, 0), (-1, -1), 2 * mm),
        ("TOPPADDING", (0, 0), (-1, -1), 2 * mm), ("BOTTOMPADDING", (0, 0), (-1, -1), 2 * mm),
    ]))
    story.extend([domain_table, PageBreak(), p("TEST TERRAIN / EXAM SPECTRUM", code), p("전 문항을 순서·사고행동·난도·배점으로 대조합니다", section)])
    spectrum_rows = [[p("번호", h3), p("단원", h3), p("사고행동", h3), p("난도", h3), p("배점", h3)]]
    for item in report["questions"]:
        spectrum_rows.append([p(item.get("number"), code), p(_domain_for_question(report, item), small), p(item.get("thinking_action"), small), p(item.get("difficulty"), small), p(item.get("points"), small)])
    spectrum_table = Table(spectrum_rows, colWidths=[18 * mm, 70 * mm, 32 * mm, 28 * mm, 32 * mm], repeatRows=1)
    spectrum_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{DEEP_INK}")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(f"#{LAB_PAPER}")]),
        ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor(f"#{MIST}")), ("ALIGN", (0, 0), (-1, -1), "CENTER"),
        ("LEFTPADDING", (0, 0), (-1, -1), 1.7 * mm), ("RIGHTPADDING", (0, 0), (-1, -1), 1.7 * mm),
        ("TOPPADDING", (0, 0), (-1, -1), 1.6 * mm), ("BOTTOMPADDING", (0, 0), (-1, -1), 1.6 * mm),
    ]))
    story.append(spectrum_table)

    ledger_groups = _balanced_chunks(_ledger_rows(report), 16)
    for page_index, group in enumerate(ledger_groups, start=1):
        story.extend([
            PageBreak(),
            p("EVIDENCE LEDGER", code),
            p(f"전 문항 증거 원장 · {page_index}/{len(ledger_groups)}", section),
        ])
        ledger = [[p("번호", h3), p("배점", h3), p("단원", h3), p("행동", h3), p("난도", h3), p("핵심 증거", h3), p("무너지는 함정", h3)]]
        for row in group:
            ledger.append([p(row["number"], code), p(row["points"], small), p(row["domain"], small), p(row["action"], small), p(row["difficulty"], small), p(row["key"], small), p(row["trap"], small)])
        ledger_table = Table(ledger, colWidths=[11 * mm, 13 * mm, 27 * mm, 15 * mm, 15 * mm, 51 * mm, 48 * mm], repeatRows=1)
        ledger_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{DEEP_INK}")), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(f"#{LAB_PAPER}")]),
            ("GRID", (0, 0), (-1, -1), 0.25, colors.HexColor(f"#{MIST}")), ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("ALIGN", (0, 0), (4, -1), "CENTER"),
            ("LEFTPADDING", (0, 0), (-1, -1), 1.2 * mm), ("RIGHTPADDING", (0, 0), (-1, -1), 1.2 * mm),
            ("TOPPADDING", (0, 0), (-1, -1), 1.4 * mm), ("BOTTOMPADDING", (0, 0), (-1, -1), 1.4 * mm),
        ]))
        story.append(ledger_table)

    for index, item in enumerate(_xray_items(report), start=1):
        story.extend([PageBreak(), p(f"QUESTION X-RAY {index}", code), p(item.get("title") or "핵심 변별 문항", section)])
        story.append(p(f"문항 {', '.join(item.get('question_numbers') or []) or '검수 필요'}", h3))
        xray = Table([
            [p("EVIDENCE", code), p(item.get("evidence") or item.get("reason"))],
            [p("BREAK BRANCHES", code), p("\n".join(f"{i + 1}. {value}" for i, value in enumerate(item.get("collapse_branches") or [item.get("collapse_point")])) )],
            [p("RECOVERY STEPS", code), p("\n".join(f"{i + 1}. {value}" for i, value in enumerate(item.get("recovery_steps") or [item.get("prescription")])) )],
            [p("LEARNING SIGNAL", code), p(item.get("learning_point") or item.get("prescription"), h3)],
        ], colWidths=[38 * mm, 142 * mm])
        xray.setStyle(TableStyle([
            ("LINEBEFORE", (0, 0), (0, -1), 3, colors.HexColor(f"#{PLASMA_BLUE}")),
            ("LINEBELOW", (0, 0), (-1, -1), 0.45, colors.HexColor(f"#{MIST}")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 3 * mm), ("RIGHTPADDING", (0, 0), (-1, -1), 3 * mm),
            ("TOPPADDING", (0, 0), (-1, -1), 3 * mm), ("BOTTOMPADDING", (0, 0), (-1, -1), 3 * mm),
        ]))
        story.append(xray)

    story.extend([PageBreak(), p("ERROR GENOME / RECOVERY PROTOCOL", code), p("오류의 원인에서 다음 행동까지", section)])
    for index, item in enumerate((report.get("failure_patterns") or [])[:4], start=1):
        genome = Table([[p(f"G{index:02d}", code), p(item.get("title"), h3), p(item.get("symptom"), small), p(item.get("cause"), small), p(item.get("prescription"), small)]], colWidths=[13 * mm, 31 * mm, 42 * mm, 42 * mm, 52 * mm])
        genome.setStyle(TableStyle([
            ("LINEBELOW", (0, 0), (-1, -1), 0.45, colors.HexColor(f"#{MIST}")), ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 2 * mm), ("RIGHTPADDING", (0, 0), (-1, -1), 2 * mm),
            ("TOPPADDING", (0, 0), (-1, -1), 2.5 * mm), ("BOTTOMPADDING", (0, 0), (-1, -1), 2.5 * mm),
        ]))
        story.append(genome)
    protocol = report.get("recovery_protocol") or {}
    story.extend([Spacer(1, 6 * mm), p("RECOVERY PROTOCOL", code)])
    protocol_table = Table([
        [p("72시간", h3), p("2주", h3), p("다음 시험", h3)],
        [
            p("\n".join(f"- {item}" for item in protocol.get("within_72_hours") or report["conclusion"].get("actions") or []), small),
            p("\n".join(f"- {item}" for item in protocol.get("within_two_weeks") or []), small),
            p("\n".join(f"- {item}" for item in protocol.get("next_exam") or []), small),
        ],
    ], colWidths=[60 * mm] * 3)
    protocol_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{LAB_PAPER}")),
        ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor(f"#{MIST}")), ("INNERGRID", (0, 0), (-1, -1), 0.35, colors.HexColor(f"#{MIST}")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 3 * mm), ("RIGHTPADDING", (0, 0), (-1, -1), 3 * mm),
        ("TOPPADDING", (0, 0), (-1, -1), 3 * mm), ("BOTTOMPADDING", (0, 0), (-1, -1), 3 * mm),
    ]))
    story.append(protocol_table)
    bands = report.get("achievement_bands") or []
    if bands:
        story.extend([Spacer(1, 6 * mm), p("ACHIEVEMENT SIGNALS", code)])
        band_rows = [[p("관측 구간", h3), p("확인 신호", h3), p("다음 처방", h3)]]
        for item in bands[:3]:
            band_rows.append([p(item.get("label"), h3), p(item.get("signal"), small), p(item.get("prescription"), small)])
        band_table = Table(band_rows, colWidths=[36 * mm, 72 * mm, 72 * mm], repeatRows=1)
        band_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(f"#{LAB_PAPER}")),
            ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor(f"#{MIST}")), ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 2.5 * mm), ("RIGHTPADDING", (0, 0), (-1, -1), 2.5 * mm),
            ("TOPPADDING", (0, 0), (-1, -1), 2 * mm), ("BOTTOMPADDING", (0, 0), (-1, -1), 2 * mm),
        ]))
        story.append(band_table)

    story.extend([PageBreak(), p("PARENT MEMO / NEXT SIGNAL", code), p("점수보다 먼저 확인할 질문", section)])
    guidance = report.get("parent_guidance") or {}
    parent_table = Table([
        [p("피할 말", h3), p("함께 확인할 질문", h3)],
        [p("\n".join(f"— {item}" for item in guidance.get("avoid") or [])), p("\n".join(f"— {item}" for item in guidance.get("recommended") or []))],
    ], colWidths=[90 * mm] * 2)
    parent_table.setStyle(TableStyle([
        ("LINEBEFORE", (0, 0), (0, -1), 3, colors.HexColor(f"#{SIGNAL_CORAL}")),
        ("LINEBEFORE", (1, 0), (1, -1), 3, colors.HexColor(f"#{PLASMA_BLUE}")),
        ("VALIGN", (0, 0), (-1, -1), "TOP"), ("LEFTPADDING", (0, 0), (-1, -1), 4 * mm), ("RIGHTPADDING", (0, 0), (-1, -1), 4 * mm),
        ("TOPPADDING", (0, 0), (-1, -1), 3 * mm), ("BOTTOMPADDING", (0, 0), (-1, -1), 3 * mm),
    ]))
    story.extend([parent_table, Spacer(1, 10 * mm), p("NEXT SIGNAL", code), p(report["conclusion"].get("headline") or summary.get("one_line"), ParagraphStyle("NextSignal", parent=section, fontSize=18, leading=25))])
    for index, action in enumerate(report["conclusion"].get("actions") or [], start=1):
        story.extend([p(f"{index:02d}  {action}", ParagraphStyle(f"Action{index}", parent=body, fontName=mono, fontSize=9.5, leading=14)), Spacer(1, 1.5 * mm)])

    doc.build(story, onFirstPage=footer, onLaterPages=footer)
    return buffer.getvalue()
