# PATH: apps/domains/tools/ppt/services.py
# PPT 생성 서비스 — python-pptx + Pillow
#
# 기능:
# - 이미지 → 슬라이드 (16:9, 4:3 지원)
# - 흑백 반전 (빔프로젝터용)
# - 배경색 설정 (검정/흰색/커스텀)
# - 이미지 비율 유지하면서 슬라이드에 최대 크기 배치

from __future__ import annotations

import io
import logging
from typing import Literal

from PIL import Image, ImageOps, ImageEnhance, ImageStat, ImageFilter
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)

# Slide dimensions in EMU (English Metric Units, 1 inch = 914400 EMU)
SLIDE_DIMENSIONS = {
    "16:9": (Inches(13.333), Inches(7.5)),
    "4:3": (Inches(10), Inches(7.5)),
}

BACKGROUND_COLORS = {
    "black": RGBColor(0, 0, 0),
    "white": RGBColor(255, 255, 255),
    "dark_gray": RGBColor(30, 30, 30),
}

# 이미지 최대 크기 (Pillow decompression bomb 방지)
MAX_IMAGE_PIXELS = 50_000_000  # 50 megapixels
# Pillow 기본 제한을 올림 (기본 178M, 충분히 안전)
Image.MAX_IMAGE_PIXELS = MAX_IMAGE_PIXELS


def validate_image_bytes(raw_bytes: bytes) -> Image.Image:
    """이미지 바이트를 검증하고 PIL Image로 반환.

    Raises:
        ValueError: 유효하지 않은 이미지인 경우.
    """
    if len(raw_bytes) < 8:
        raise ValueError("Too small to be a valid image")

    try:
        img = Image.open(io.BytesIO(raw_bytes))
        img.verify()  # 구조 검증 (verify 후 재로딩 필요)
    except Exception as exc:
        raise ValueError(f"Invalid image data: {exc}") from exc

    # verify() 후 재로딩 (verify는 이미지를 소비함)
    img = Image.open(io.BytesIO(raw_bytes))

    # 크기 검증
    w, h = img.size
    if w <= 0 or h <= 0:
        raise ValueError(f"Invalid image dimensions: {w}x{h}")
    if w * h > MAX_IMAGE_PIXELS:
        raise ValueError(f"Image too large: {w}x{h} ({w * h} pixels)")

    return img


def _auto_enhance_document(img: Image.Image) -> Image.Image:
    """문서용 자동 보정 — "읽기 쉬운 흑백" 목표.

    목적: 빔프로젝터에서 문서/시험지가 선명하게 보이도록.
    원칙:
    - "밝게"가 아니라 "글자 선명, 배경 깨끗"
    - 배경 watermark/로고가 과도하게 밝아지면 안 됨
    - histogram 기반 대비 정규화 (단순 brightness 증가 아님)
    - 선명도 보강으로 글자 edge 강화
    """
    gray = img.convert("L")
    stat = ImageStat.Stat(gray)
    mean_val = stat.mean[0]
    stddev = stat.stddev[0] if stat.stddev else 0

    # 이미 고대비 흑백 (mean~128, stddev 높음) → 보정 최소화
    if stddev > 80 and 80 < mean_val < 180:
        # 선명도만 살짝 보강
        img = ImageEnhance.Sharpness(img).enhance(1.3)
        return img

    # --- histogram 기반 대비 정규화 (autocontrast) ---
    # cutoff: 상하위 0.5% 클리핑 → 연한 배경/워터마크는 날리지 않으면서
    # 글자와 배경의 대비를 자연스럽게 확장
    if img.mode == "RGB":
        img = ImageOps.autocontrast(img, cutoff=0.5)
    elif img.mode == "L":
        img = ImageOps.autocontrast(img, cutoff=0.5)
        img = img.convert("RGB")

    # --- 대비 보강 (제한적) ---
    # stddev가 낮으면 = 대비 부족 (흐릿한 스캔 등) → 대비만 올림
    if stddev < 40:
        img = ImageEnhance.Contrast(img).enhance(1.4)
    elif stddev < 60:
        img = ImageEnhance.Contrast(img).enhance(1.2)

    # --- 선명도 보강 (글자 edge 강화) ---
    img = ImageEnhance.Sharpness(img).enhance(1.5)

    return img


def _needs_processing(
    invert: bool, grayscale: bool, auto_enhance: bool,
    brightness: float, contrast: float,
) -> bool:
    """효과가 하나라도 적용되는지 확인."""
    return invert or grayscale or auto_enhance or brightness != 1.0 or contrast != 1.0


def _detect_format(image_bytes: bytes) -> str:
    """이미지 바이트에서 원본 포맷 감지."""
    if image_bytes[:3] == b"\xff\xd8\xff":
        return "JPEG"
    if image_bytes[:8] == b"\x89PNG\r\n\x1a\n":
        return "PNG"
    return "PNG"  # 기타 → PNG


def _process_image(
    image_bytes: bytes,
    *,
    invert: bool = False,
    grayscale: bool = False,
    auto_enhance: bool = False,
    brightness: float = 1.0,
    contrast: float = 1.0,
) -> bytes:
    """이미지 전처리: 반전, 그레이스케일, 밝기/대비 조절 등.

    효과가 없으면 원본 바이트를 그대로 반환하여 화질 손실을 방지한다.
    JPEG 원본은 JPEG로 유지 (quality=95), PNG 원본은 PNG로 유지.

    Args:
        invert: 흑백 반전.
        grayscale: 그레이스케일 변환.
        auto_enhance: 자동 밝기/대비 보정 (빔프로젝터용).
        brightness: 밝기 배수 (1.0=원본).
        contrast: 대비 배수 (1.0=원본).

    Raises:
        ValueError: 유효하지 않은 이미지.
    """
    # 유효성 검증은 항상 수행
    validate_image_bytes(image_bytes)

    # 효과 없으면 원본 그대로 반환 (화질 보존)
    if not _needs_processing(invert, grayscale, auto_enhance, brightness, contrast):
        return image_bytes

    orig_format = _detect_format(image_bytes)
    img = Image.open(io.BytesIO(image_bytes))

    # EXIF 회전 보정
    img = ImageOps.exif_transpose(img)

    # RGBA → RGB 변환 (PPT 호환)
    if img.mode == "RGBA":
        background = Image.new("RGB", img.size, (255, 255, 255))
        background.paste(img, mask=img.split()[3])
        img = background
    elif img.mode not in ("RGB", "L"):
        img = img.convert("RGB")

    if img.mode == "L":
        img = img.convert("RGB")

    if grayscale:
        img = ImageOps.grayscale(img)
        img = img.convert("RGB")

    if invert:
        img = ImageOps.invert(img)

    # 자동 보정 (문서용: 선명한 흑백 출력 목표)
    if auto_enhance:
        img = _auto_enhance_document(img)

    # 수동 밝기/대비 조절 (자동 보정 위에 추가 적용 가능)
    if brightness != 1.0:
        brightness = max(0.2, min(3.0, brightness))
        img = ImageEnhance.Brightness(img).enhance(brightness)
        # 밝기를 올리면 대비가 떨어지므로 자동 보상
        if brightness > 1.0 and contrast == 1.0:
            compensation = 1.0 + (brightness - 1.0) * 0.3
            img = ImageEnhance.Contrast(img).enhance(compensation)

    if contrast != 1.0:
        contrast = max(0.2, min(3.0, contrast))
        img = ImageEnhance.Contrast(img).enhance(contrast)

    # 원본 포맷 유지하여 저장 (화질 보존)
    out = io.BytesIO()
    if orig_format == "JPEG":
        img.save(out, format="JPEG", quality=95, subsampling=0)
    else:
        img.save(out, format="PNG", optimize=True)
    out.seek(0)
    return out.read()


def _fit_image_to_slide(
    img_width: int,
    img_height: int,
    slide_width: int,
    slide_height: int,
    fit_mode: str = "contain",
) -> tuple[int, int, int, int]:
    """이미지를 슬라이드에 맞추는 위치/크기 계산 (EMU 단위).

    fit_mode:
      - "contain": 비율 유지, 슬라이드 안에 최대 크기 (기본)
      - "cover": 비율 유지, 슬라이드 전체 덮기 (일부 잘림)
      - "stretch": 비율 무시, 슬라이드 꽉 채움

    Returns: (left, top, width, height) in EMU.
    """
    if fit_mode == "stretch":
        return 0, 0, slide_width, slide_height

    # ZeroDivisionError 방어
    if img_width <= 0 or img_height <= 0:
        return 0, 0, slide_width, slide_height

    img_ratio = img_width / img_height
    slide_ratio = slide_width / slide_height

    if fit_mode == "cover":
        if img_ratio > slide_ratio:
            height = slide_height
            width = int(height * img_ratio)
        else:
            width = slide_width
            height = int(width / img_ratio)
    else:  # contain
        if img_ratio > slide_ratio:
            width = slide_width
            height = int(width / img_ratio)
        else:
            height = slide_height
            width = int(height * img_ratio)

    left = (slide_width - width) // 2
    top = (slide_height - height) // 2
    return left, top, width, height


def _get_blank_layout(prs: Presentation):
    """빈 슬라이드 레이아웃을 안전하게 찾는다.

    index 6이 항상 blank라는 보장이 없으므로 이름으로 검색 후 fallback.
    """
    # 이름으로 검색 (Blank, 빈 슬라이드 등)
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if name in ("blank", "빈 슬라이드", "빈 화면"):
            return layout

    # placeholder가 가장 적은 레이아웃 선택
    min_ph = None
    best = None
    for layout in prs.slide_layouts:
        ph_count = len(layout.placeholders)
        if min_ph is None or ph_count < min_ph:
            min_ph = ph_count
            best = layout

    if best is not None:
        return best

    # 최후 fallback: 마지막 레이아웃 (보통 blank)
    return prs.slide_layouts[-1]


def generate_ppt(
    images: list[tuple[str, bytes]],
    *,
    aspect_ratio: Literal["16:9", "4:3"] = "16:9",
    background: str = "black",
    fit_mode: str = "contain",
    invert: bool = False,
    grayscale: bool = False,
    auto_enhance: bool = False,
    brightness: float = 1.0,
    contrast: float = 1.0,
    per_slide_settings: list[dict] | None = None,
) -> bytes:
    """PPT 생성.

    Args:
        images: (파일명, 이미지바이트) 리스트. 순서대로 슬라이드 생성.
        aspect_ratio: 슬라이드 비율 ("16:9" 또는 "4:3").
        background: 배경색 ("black", "white", "dark_gray" 또는 hex "#RRGGBB").
        fit_mode: 이미지 배치 모드 ("contain", "cover", "stretch").
        invert: True면 모든 이미지 흑백 반전.
        grayscale: True면 모든 이미지 그레이스케일 변환.
        auto_enhance: True면 어두운 이미지 자동 밝기/대비 보정.
        brightness: 밝기 배수 (1.0=원본).
        contrast: 대비 배수 (1.0=원본).
        per_slide_settings: 슬라이드별 개별 설정.

    Returns:
        PPTX 파일 바이트.

    Raises:
        ValueError: 이미지가 비어있거나 유효하지 않은 경우.
    """
    if not images:
        raise ValueError("No images provided")

    prs = Presentation()

    # 슬라이드 크기 설정
    slide_w, slide_h = SLIDE_DIMENSIONS.get(aspect_ratio, SLIDE_DIMENSIONS["16:9"])
    prs.slide_width = slide_w
    prs.slide_height = slide_h

    # 배경색 결정
    bg_color = BACKGROUND_COLORS.get(background)
    if bg_color is None and background.startswith("#") and len(background) == 7:
        try:
            r = int(background[1:3], 16)
            g = int(background[3:5], 16)
            b = int(background[5:7], 16)
            bg_color = RGBColor(r, g, b)
        except ValueError:
            bg_color = BACKGROUND_COLORS["black"]
    elif bg_color is None:
        bg_color = BACKGROUND_COLORS["black"]

    blank_layout = _get_blank_layout(prs)

    failed_slides = []
    for idx, (filename, raw_bytes) in enumerate(images):
        # 슬라이드별 개별 설정
        slide_invert = invert
        slide_grayscale = grayscale
        slide_auto_enhance = auto_enhance
        slide_brightness = brightness
        slide_contrast = contrast
        if per_slide_settings and idx < len(per_slide_settings):
            ss = per_slide_settings[idx]
            if isinstance(ss, dict):
                slide_invert = ss.get("invert", invert)
                slide_grayscale = ss.get("grayscale", grayscale)
                slide_auto_enhance = ss.get("auto_enhance", auto_enhance)
                slide_brightness = float(ss.get("brightness", brightness))
                slide_contrast = float(ss.get("contrast", contrast))

        # 이미지 전처리
        try:
            processed = _process_image(
                raw_bytes,
                invert=slide_invert,
                grayscale=slide_grayscale,
                auto_enhance=slide_auto_enhance,
                brightness=slide_brightness,
                contrast=slide_contrast,
            )
        except (ValueError, Exception) as exc:
            logger.warning("이미지 처리 실패: %s (idx=%d): %s", filename, idx, exc)
            failed_slides.append(idx)
            continue  # 실패한 이미지는 건너뜀 (원본 삽입 시도 대신 안전하게 skip)

        # 이미지 크기 파악
        img = Image.open(io.BytesIO(processed))
        img_w, img_h = img.size

        # 슬라이드 추가
        slide = prs.slides.add_slide(blank_layout)

        # 배경색 설정
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # 이미지 위치/크기 계산 (EMU)
        left, top, width, height = _fit_image_to_slide(
            img_w, img_h,
            prs.slide_width, prs.slide_height,
            fit_mode=fit_mode,
        )

        # 이미지 추가
        img_stream = io.BytesIO(processed)
        slide.shapes.add_picture(img_stream, left, top, width, height)

    # 유효한 슬라이드가 하나도 없으면 에러
    successful = len(images) - len(failed_slides)
    if successful == 0:
        raise ValueError(f"All {len(images)} images failed to process")

    if failed_slides:
        logger.info("PPT 생성: %d/%d 슬라이드 성공, 실패=%s", successful, len(images), failed_slides)

    # PPTX 바이트 출력
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
