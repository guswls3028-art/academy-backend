# PATH: apps/domains/tools/ppt/tests.py
# PPT 생성 서비스 테스트 — 하드닝 검증

import io
import time
import zipfile

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from .services import generate_ppt, validate_image_bytes, _process_image, _fit_image_to_slide


def make_img(w=100, h=100, mode="RGB", fmt="PNG"):
    color = (100, 150, 200) if mode == "RGB" else (100, 150, 200, 128)
    img = Image.new(mode, (w, h), color=color)
    buf = io.BytesIO()
    img.save(buf, format=fmt)
    return buf.getvalue()


def run_all():
    passed = 0
    failed = 0

    def test(name, fn):
        nonlocal passed, failed
        try:
            fn()
            print(f"  PASS: {name}")
            passed += 1
        except AssertionError as e:
            print(f"  FAIL: {name} - {e}")
            failed += 1
        except Exception as e:
            print(f"  FAIL: {name} - {type(e).__name__}: {e}")
            failed += 1

    # === 1. Happy Path ===
    print("=== 1. Happy Path ===")

    def t_basic():
        r = generate_ppt([("a.png", make_img())])
        assert len(r) > 100

    def t_4_3():
        r = generate_ppt([("a.png", make_img())], aspect_ratio="4:3", background="white")
        assert len(r) > 100

    def t_invert_gray():
        r = generate_ppt([("a.png", make_img())], invert=True, grayscale=True)
        assert len(r) > 100

    def t_cover():
        r = generate_ppt([("a.png", make_img(800, 400))], fit_mode="cover")
        assert len(r) > 100

    def t_stretch():
        r = generate_ppt([("a.png", make_img())], fit_mode="stretch")
        assert len(r) > 100

    def t_multi():
        imgs = [("a.png", make_img()), ("b.png", make_img(200, 100)), ("c.png", make_img(50, 300))]
        r = generate_ppt(imgs)
        assert len(r) > 100

    def t_hex_bg():
        r = generate_ppt([("a.png", make_img())], background="#ff00aa")
        assert len(r) > 100

    def t_per_slide():
        r = generate_ppt(
            [("a.png", make_img()), ("b.png", make_img())],
            per_slide_settings=[{"invert": True}, {"grayscale": True}],
        )
        assert len(r) > 100

    test("basic 16:9 black", t_basic)
    test("4:3 white bg", t_4_3)
    test("invert + grayscale", t_invert_gray)
    test("cover fit", t_cover)
    test("stretch fit", t_stretch)
    test("multiple slides", t_multi)
    test("hex background", t_hex_bg)
    test("per_slide settings", t_per_slide)

    # === 2. Invalid Inputs ===
    print("\n=== 2. Invalid Inputs ===")

    def t_empty():
        try:
            generate_ppt([])
            raise AssertionError("should have raised ValueError")
        except ValueError:
            pass

    def t_invalid_bytes():
        try:
            validate_image_bytes(b"not an image at all garbage data")
            raise AssertionError("should have raised ValueError")
        except ValueError:
            pass

    def t_too_small():
        try:
            validate_image_bytes(b"abc")
            raise AssertionError("should have raised ValueError")
        except ValueError as e:
            assert "Too small" in str(e)

    test("empty images raises ValueError", t_empty)
    test("invalid image bytes raises ValueError", t_invalid_bytes)
    test("too small bytes raises ValueError", t_too_small)

    # === 3. Edge Cases ===
    print("\n=== 3. Edge Cases ===")

    def t_rgba():
        r = generate_ppt([("a.png", make_img(mode="RGBA"))])
        assert len(r) > 100

    def t_panorama():
        r = generate_ppt([("a.png", make_img(2000, 100))])
        assert len(r) > 100

    def t_tall():
        r = generate_ppt([("a.png", make_img(100, 2000))])
        assert len(r) > 100

    def t_1x1():
        r = generate_ppt([("a.png", make_img(1, 1))])
        assert len(r) > 100

    def t_bad_bg():
        r = generate_ppt([("a.png", make_img())], background="notacolor")
        assert len(r) > 100

    def t_bad_hex():
        r = generate_ppt([("a.png", make_img())], background="#xyz")
        assert len(r) > 100

    def t_per_slide_nondict():
        r = generate_ppt([("a.png", make_img())], per_slide_settings=["invalid"])
        assert len(r) > 100

    def t_jpeg():
        r = generate_ppt([("a.jpg", make_img(fmt="JPEG"))])
        assert len(r) > 100

    def t_bmp():
        r = generate_ppt([("a.bmp", make_img(fmt="BMP"))])
        assert len(r) > 100

    test("RGBA image", t_rgba)
    test("very wide image (panorama)", t_panorama)
    test("very tall image", t_tall)
    test("1x1 pixel image", t_1x1)
    test("invalid bg falls back to black", t_bad_bg)
    test("invalid hex bg falls back to black", t_bad_hex)
    test("per_slide with non-dict entry", t_per_slide_nondict)
    test("JPEG format", t_jpeg)
    test("BMP format", t_bmp)

    # === 4. ZeroDivisionError defense ===
    print("\n=== 4. ZeroDivisionError defense ===")

    def t_zero_h():
        result = _fit_image_to_slide(100, 0, 1000, 750)
        assert result == (0, 0, 1000, 750), f"got {result}"

    def t_zero_w():
        result = _fit_image_to_slide(0, 100, 1000, 750)
        assert result == (0, 0, 1000, 750), f"got {result}"

    test("_fit_image zero height", t_zero_h)
    test("_fit_image zero width", t_zero_w)

    # === 5. Output Validity ===
    print("\n=== 5. Output Validity ===")

    def t_valid_zip():
        r = generate_ppt([("a.png", make_img())])
        z = zipfile.ZipFile(io.BytesIO(r))
        names = z.namelist()
        assert "[Content_Types].xml" in names, f"missing: {names[:5]}"
        assert any("slide1.xml" in n for n in names), f"no slide1: {names[:10]}"
        z.close()

    def t_opens_as_prs():
        r = generate_ppt([("a.png", make_img()), ("b.png", make_img())])
        prs = Presentation(io.BytesIO(r))
        assert len(prs.slides) == 2, f"expected 2, got {len(prs.slides)}"

    def t_dims_16_9():
        r = generate_ppt([("a.png", make_img())], aspect_ratio="16:9")
        prs = Presentation(io.BytesIO(r))
        assert prs.slide_width == Inches(13.333), f"width: {prs.slide_width}"

    def t_dims_4_3():
        r = generate_ppt([("a.png", make_img())], aspect_ratio="4:3")
        prs = Presentation(io.BytesIO(r))
        assert prs.slide_width == Inches(10), f"width: {prs.slide_width}"

    def t_slide_has_image():
        r = generate_ppt([("a.png", make_img())])
        prs = Presentation(io.BytesIO(r))
        slide = prs.slides[0]
        shapes = list(slide.shapes)
        assert len(shapes) >= 1, f"no shapes on slide"
        # At least one shape should be a picture
        assert any(hasattr(s, "image") for s in shapes), "no image shape found"

    test("output is valid PPTX (zip)", t_valid_zip)
    test("output opens as Presentation", t_opens_as_prs)
    test("slide dimensions 16:9", t_dims_16_9)
    test("slide dimensions 4:3", t_dims_4_3)
    test("slide contains image shape", t_slide_has_image)

    # === 6. Mixed valid+invalid images ===
    print("\n=== 6. Mixed valid+invalid images ===")

    def t_skip_bad():
        r = generate_ppt([
            ("good.png", make_img()),
            ("bad.bin", b"not-an-image-at-all-xyzxyz" * 10),
            ("good2.png", make_img()),
        ])
        prs = Presentation(io.BytesIO(r))
        assert len(prs.slides) == 2, f"expected 2 (1 skipped), got {len(prs.slides)}"

    def t_all_bad():
        try:
            generate_ppt([("bad1.bin", b"xxxx" * 100), ("bad2.bin", b"yyyy" * 100)])
            raise AssertionError("should have raised ValueError")
        except ValueError as e:
            assert "All" in str(e)

    test("bad image among good ones is skipped", t_skip_bad)
    test("all bad images raises ValueError", t_all_bad)

    # === 7. Stress ===
    print("\n=== 7. Stress ===")

    def t_50_slides():
        t = time.time()
        imgs = [(f"{i}.png", make_img(400, 300)) for i in range(50)]
        r = generate_ppt(imgs)
        elapsed = time.time() - t
        prs = Presentation(io.BytesIO(r))
        print(f"    50 slides: {len(r)} bytes in {elapsed:.1f}s")
        assert len(prs.slides) == 50, f"expected 50, got {len(prs.slides)}"
        assert elapsed < 60, f"too slow: {elapsed:.1f}s"

    def t_large_images():
        t = time.time()
        imgs = [(f"{i}.png", make_img(2000, 1500)) for i in range(10)]
        r = generate_ppt(imgs, invert=True, grayscale=True)
        elapsed = time.time() - t
        print(f"    10 large slides (2000x1500 + effects): {len(r)} bytes in {elapsed:.1f}s")
        assert elapsed < 60, f"too slow: {elapsed:.1f}s"

    test("50 slides performance", t_50_slides)
    test("10 large images with effects", t_large_images)

    # === 8. Process Image Edge Cases ===
    print("\n=== 8. Process Image Edge Cases ===")

    def t_process_invert():
        r = _process_image(make_img(), invert=True)
        img = Image.open(io.BytesIO(r))
        assert img.mode == "RGB"

    def t_process_grayscale():
        r = _process_image(make_img(), grayscale=True)
        img = Image.open(io.BytesIO(r))
        assert img.mode == "RGB"

    def t_process_both():
        r = _process_image(make_img(), invert=True, grayscale=True)
        img = Image.open(io.BytesIO(r))
        assert img.mode == "RGB"

    def t_process_rgba():
        # 효과 없으면 원본 유지, 효과 있으면 RGB 변환
        rgba_bytes = make_img(mode="RGBA")
        r = _process_image(rgba_bytes)
        assert r == rgba_bytes, "No effects = original preserved"
        r2 = _process_image(rgba_bytes, invert=True)
        img = Image.open(io.BytesIO(r2))
        assert img.mode == "RGB", "With effects = converted to RGB"

    test("_process_image invert", t_process_invert)
    test("_process_image grayscale", t_process_grayscale)
    test("_process_image invert+grayscale", t_process_both)
    test("_process_image RGBA→RGB", t_process_rgba)

    # === Summary ===
    print(f"\n{'='*40}")
    print(f"Results: {passed} passed, {failed} failed")
    if failed:
        print("SOME TESTS FAILED")
        return 1
    else:
        print("ALL TESTS PASSED")
        return 0


if __name__ == "__main__":
    import os, sys
    os.environ.setdefault("DJANGO_SETTINGS_MODULE", "apps.api.config.settings.dev")
    import django
    django.setup()
    sys.exit(run_all())
